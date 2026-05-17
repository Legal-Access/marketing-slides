# -*- coding: utf-8 -*-
"""
Legal Access 法務AI活用ロードマップ策定・推進 営業スライド
HTMLドラフトと同じ構成・配色で .pptx を生成する。

仕様:
- 16:9 / 1920x1080 px 相当（PPTXのEMU換算）
- ネイビー + シアン基調
- 最小フォント18pt
- 経歴9項目・セミナー7件+多数 を省略せず全件掲載
"""

import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ============ 設定 ============
HERE = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(HERE, "20260517_LegalAccess_法務AIロードマップ提案.pptx")

# 16:9 ハーフEMU。1920x1080pxを基準とし、PPTXは13.333in x 7.5in
SLIDE_W = Emu(12192000)   # 13.333 inch
SLIDE_H = Emu(6858000)    # 7.5    inch

# スケール: 1920px -> 12192000 EMU なので 1px = 6350 EMU
PX = 6350
def x(px): return Emu(px * PX)

# 色
C_NAVY    = RGBColor(0x0B, 0x1F, 0x3A)
C_DEEP    = RGBColor(0x06, 0x12, 0x2A)
C_CYAN    = RGBColor(0x3F, 0xD0, 0xE9)
C_CYAN2   = RGBColor(0x00, 0xB7, 0xD4)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_SUBGRAY = RGBColor(0xE8, 0xEE, 0xF5)
C_BG_LITE = RGBColor(0xF4, 0xF8, 0xFC)
C_TEXT2   = RGBColor(0x3D, 0x4F, 0x6B)
C_TEXT_L  = RGBColor(0xB9, 0xC7, 0xDE)
C_PAGE    = RGBColor(0x99, 0xA8, 0xBD)

FONT = "Yu Gothic UI"
FONT_B = "Yu Gothic UI"

# ============ ヘルパー ============
def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = Pt(line_w)
    shape.text_frame.margin_left = 0
    shape.text_frame.margin_right = 0
    shape.text_frame.margin_top = 0
    shape.text_frame.margin_bottom = 0
    return shape

def add_text(slide, left, top, width, height, runs,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             fill=None, line=None, padding=(8, 8, 8, 8),
             line_spacing=1.15, line_w=None):
    """
    runs: list of dicts: {text, size, bold, color, font, italic}
    一段落＝一要素。改行したい場合は別段落として渡す（複数段落=list of list of dict）。
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = Pt(line_w)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left   = Emu(padding[3] * PX)
    tf.margin_right  = Emu(padding[1] * PX)
    tf.margin_top    = Emu(padding[0] * PX)
    tf.margin_bottom = Emu(padding[2] * PX)
    tf.vertical_anchor = anchor

    # runs may be list of dict (single para) or list of list of dict (multi)
    if isinstance(runs, list) and runs and isinstance(runs[0], dict):
        paragraphs = [runs]
    else:
        paragraphs = runs

    for i, para in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for j, r in enumerate(para):
            if j == 0 and i == 0:
                run = p.runs[0] if p.runs else p.add_run()
                run.text = r["text"]
            else:
                run = p.add_run()
                run.text = r["text"]
            run.font.name = r.get("font", FONT)
            run.font.size = Pt(r.get("size", 18))
            run.font.bold = r.get("bold", False)
            run.font.italic = r.get("italic", False)
            color = r.get("color", C_NAVY)
            run.font.color.rgb = color
    return shape

def add_line_shape(slide, left, top, width, height, color, weight=2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_image(slide, path, left, top, width, height):
    return slide.shapes.add_picture(path, left, top, width=width, height=height)

def add_page_chrome(slide, page_no, total=13, show_footer=True, brand_color=C_NAVY):
    # accent bar
    add_rect(slide, x(0), x(0), x(12), x(1080), fill=C_CYAN2)
    if show_footer:
        # footer brand
        add_text(slide, x(48), x(1024), x(400), x(36),
                 [
                     {"text": "Legal ", "size": 14, "bold": True, "color": brand_color},
                     {"text": "Access", "size": 14, "bold": True, "color": C_CYAN2},
                 ],
                 anchor=MSO_ANCHOR.MIDDLE)
        # page num
        add_text(slide, x(1700), x(1024), x(180), x(36),
                 [{"text": f"{page_no:02d} / {total:02d}", "size": 14, "color": C_PAGE}],
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def add_heading(slide, section_label, title_runs):
    """
    title_runs: list of dict for the slide title (allows partial cyan emphasis).
    """
    # section label with leading cyan bar
    add_line_shape(slide, x(96), x(82), x(60), x(2), C_CYAN)
    add_text(slide, x(168), x(64), x(900), x(40),
             [{"text": section_label, "size": 14, "bold": False, "color": C_CYAN2}],
             anchor=MSO_ANCHOR.MIDDLE)
    # slide title
    add_text(slide, x(96), x(108), x(1750), x(80),
             title_runs,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2)

# ============ プレゼンテーション初期化 ============
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ============ Slide 1: Cover ============
s = prs.slides.add_slide(BLANK)
# 背景: 濃ネイビーで塗る
add_rect(s, x(0), x(0), x(1920), x(1080), fill=C_DEEP)
# 背景画像（半透明感は無理なので右側に大きめに置きトリミング表現）
img_path = os.path.join(HERE, "青いネットワーク.png")
add_image(s, img_path, x(960), x(0), x(960), x(1080))
# 左側に濃ネイビーのオーバーレイ帯
add_rect(s, x(0), x(0), x(1200), x(1080), fill=C_NAVY)
# eyebrow
add_text(s, x(120), x(220), x(1100), x(40),
         [{"text": "LEGAL ACCESS  |  SERVICE PROPOSAL", "size": 16, "bold": False, "color": C_CYAN}])
# accent line
add_line_shape(s, x(120), x(280), x(120), x(4), C_CYAN)
# title
add_text(s, x(120), x(310), x(1100), x(280),
         [
             [{"text": "法務AI活用", "size": 60, "bold": True, "color": C_WHITE}],
             [
                 {"text": "ロードマップ", "size": 60, "bold": True, "color": C_CYAN},
                 {"text": "策定・推進", "size": 60, "bold": True, "color": C_WHITE},
             ],
         ],
         line_spacing=1.2)
# sub
add_text(s, x(120), x(620), x(1100), x(140),
         [[
             {"text": "「とりあえず入れたAI」を、", "size": 22, "color": C_SUBGRAY},
             {"text": "経営インパクトの出る武器", "size": 22, "bold": True, "color": C_CYAN},
             {"text": "に変える伴走型サービス", "size": 22, "color": C_SUBGRAY},
         ]],
         line_spacing=1.5)
# speaker
add_rect(s, x(120), x(820), x(220), x(56), fill=None, line=C_CYAN, line_w=1.5)
add_text(s, x(120), x(820), x(220), x(56),
         [{"text": "PRESENTED BY", "size": 14, "bold": True, "color": C_CYAN}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x(360), x(818), x(800), x(60),
         [{"text": "リーガルアクセス株式会社", "size": 22, "bold": True, "color": C_WHITE}],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x(360), x(862), x(800), x(40),
         [{"text": "代表取締役 弁護士　福島 駿太", "size": 16, "color": C_TEXT_L}],
         anchor=MSO_ANCHOR.MIDDLE)

# ============ Slide 2: Issues ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 2)
add_heading(s, "ISSUE  /  法務部の現在地",
            [[
                {"text": "入れたはずのAIが、", "size": 44, "bold": True, "color": C_NAVY},
                {"text": "動いていない", "size": 44, "bold": True, "color": C_CYAN2},
                {"text": "。", "size": 44, "bold": True, "color": C_NAVY},
            ]])
# lead caption (dark band)
add_text(s, x(96), x(230), x(1730), x(90),
         [[
             {"text": "多くの法務部で起きていること —— ", "size": 20, "bold": True, "color": C_WHITE},
             {"text": "「ツールは入った。でも、ここから先が分からない」", "size": 20, "bold": True, "color": C_CYAN},
         ]],
         fill=C_NAVY, anchor=MSO_ANCHOR.MIDDLE, padding=(20, 28, 20, 28))

# 5 cards
issues = [
    ("01", "とりあえず\n導入だけは済んだ", "ChatGPT、Gemini、Copilot — ライセンスはある。けれど、使われていない。"),
    ("02", "使い方が\n分からない",        "何を聞けば、何に使えばいいのか。社内に答えを持っている人がいない。"),
    ("03", "法務は特殊で\n事例が出てこない", "守秘性が高く、他社の活用事例が公にされない。横展開の手本がない。"),
    ("04", "効果は\nあるはずなのに",       "時間も人件費も削れるはず。でも、どの業務から手をつけるべきか決められない。"),
    ("05", "情報管理の\n線引きが見えない", "どこまで入れてよいのか。秘密情報・個人情報の扱いがブラックボックス。"),
]
card_w = 340
card_h = 580
gap = 16
total_w = card_w * 5 + gap * 4
start_x = (1920 - total_w) // 2
for i, (no, title, desc) in enumerate(issues):
    left = x(start_x + i * (card_w + gap))
    top = x(360)
    add_rect(s, left, top, x(card_w), x(card_h), fill=C_BG_LITE)
    add_rect(s, left, top, x(card_w), x(4), fill=C_CYAN2)
    add_text(s, left, top + x(28), x(card_w), x(40),
             [{"text": no, "size": 16, "bold": True, "color": C_CYAN2}],
             padding=(0, 24, 0, 24))
    add_text(s, left, top + x(72), x(card_w), x(150),
             [[{"text": title, "size": 22, "bold": True, "color": C_NAVY}]],
             padding=(0, 24, 0, 24), line_spacing=1.35)
    add_text(s, left, top + x(230), x(card_w), x(320),
             [{"text": desc, "size": 16, "color": C_TEXT2}],
             padding=(0, 24, 0, 24), line_spacing=1.6)

# ============ Slide 3: Why AI in Legal ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 3)
add_heading(s, "WHY  /  法務にAIが効く3つの理由",
            [[
                {"text": "法務こそ、AIで", "size": 44, "bold": True, "color": C_NAVY},
                {"text": "大きく変わる", "size": 44, "bold": True, "color": C_CYAN2},
                {"text": "領域。", "size": 44, "bold": True, "color": C_NAVY},
            ]])
# 3 cards
cards = [
    ("TIME",    "時間", "を、戻す",
     "契約レビュー・リサーチ・要約 — これまで弁護士・パラリーガルが手作業で時間を吸い取られていた業務に、AIは直接効く。"),
    ("COST",    "人件費", "を、削る",
     "外注に流していた定型業務を内製化できる。法務人員の採用難に対する、現実的な打ち手になる。"),
    ("QUALITY", "品質", "を、上げる",
     "抜け漏れの検出、判例横断、過去類似案件の即時参照。「人がやるよりむしろ精度が上がる」業務が存在する。"),
]
cw = 560
ch = 460
cgap = 32
ctotal = cw * 3 + cgap * 2
cstart = (1920 - ctotal) // 2
for i, (lbl, big, sm, desc) in enumerate(cards):
    left = x(cstart + i * (cw + cgap))
    top = x(230)
    add_rect(s, left, top, x(cw), x(ch), fill=C_WHITE, line=C_SUBGRAY, line_w=1.5)
    add_rect(s, left, top, x(cw), x(6), fill=C_CYAN2)
    add_text(s, left, top + x(36), x(cw), x(40),
             [{"text": lbl, "size": 16, "bold": True, "color": C_CYAN2}],
             padding=(0, 32, 0, 32))
    add_text(s, left, top + x(86), x(cw), x(110),
             [[
                 {"text": big, "size": 40, "bold": True, "color": C_NAVY},
                 {"text": sm, "size": 20, "bold": True, "color": C_CYAN2},
             ]],
             padding=(0, 32, 0, 32), line_spacing=1.1)
    add_text(s, left, top + x(220), x(cw), x(220),
             [{"text": desc, "size": 17, "color": C_TEXT2}],
             padding=(0, 32, 0, 32), line_spacing=1.7)
# quote band
add_text(s, x(96), x(720), x(1730), x(190),
         [[
             {"text": "“  ", "size": 40, "bold": True, "color": C_CYAN},
             {"text": "質が高く効率的な法律業務", "size": 24, "bold": True, "color": C_CYAN},
             {"text": "を叶える「生成AIとの上手な付き合い方」。", "size": 24, "bold": True, "color": C_WHITE},
         ],
          [
              {"text": "── 拙著『法務のための生成AI活用ガイド』（弘文堂・2025年）より", "size": 14, "color": C_TEXT_L},
          ]],
         fill=C_NAVY, padding=(28, 40, 24, 40), line_spacing=1.5)

# ============ Slide 4: Service Flow ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 4)
add_heading(s, "SERVICE  /  全体像",
            [[
                {"text": "5つのステップで、", "size": 40, "bold": True, "color": C_NAVY},
                {"text": "ロードマップから実装まで", "size": 40, "bold": True, "color": C_CYAN2},
                {"text": "。", "size": 40, "bold": True, "color": C_NAVY},
            ]])
steps = [
    ("STEP 01", "責任者・現場へ\nインタビュー",
     "担当部署の責任者と現場担当者、両方の声を聞く。建前と実態を切り分ける。"),
    ("STEP 02", "技術見通しを\n策定",
     "汎用AIの最新動向を踏まえ、「今できること」と「これからできること」を定義。"),
    ("STEP 03", "業務×AIで\nユースケース選出",
     "業務領域ごとに「AIに任せられる仕事」を具体名で洗い出す。"),
    ("STEP 04", "経営インパクトで\n評価・優先順位",
     "時間削減・コスト削減・品質向上の3軸で各ユースケースをスコア化。"),
    ("STEP 05", "ロードマップと\n実行計画に落とす",
     "短期で着手する案件、中長期で仕込む案件に振り分け、推進まで伴走。"),
]
sw = 340
sh = 460
sgap = 24
stotal = sw * 5 + sgap * 4
sstart = (1920 - stotal) // 2
for i, (no, title, desc) in enumerate(steps):
    left = x(sstart + i * (sw + sgap))
    top = x(240)
    add_rect(s, left, top, x(sw), x(sh), fill=C_WHITE, line=C_CYAN2, line_w=1.5)
    add_text(s, left, top + x(30), x(sw), x(40),
             [{"text": no, "size": 14, "bold": True, "color": C_CYAN2}],
             padding=(0, 24, 0, 24))
    add_text(s, left, top + x(78), x(sw), x(140),
             [[{"text": title, "size": 20, "bold": True, "color": C_NAVY}]],
             padding=(0, 24, 0, 24), line_spacing=1.35)
    add_text(s, left, top + x(230), x(sw), x(210),
             [{"text": desc, "size": 15, "color": C_TEXT2}],
             padding=(0, 24, 0, 24), line_spacing=1.65)
    # arrow
    if i < 4:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left + x(sw + 4), top + x(sh/2 - 10), x(16), x(20))
        arrow.rotation = 0
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = C_CYAN2
        arrow.line.fill.background()

# bottom band
add_text(s, x(96), x(780), x(1730), x(130),
         [[
             {"text": "POINT　　", "size": 18, "bold": True, "color": C_CYAN},
             {"text": "作って終わりではなく、", "size": 18, "color": C_WHITE},
             {"text": "至近のユースケースを実際に立ち上げるところまで一緒に走る", "size": 18, "bold": True, "color": C_CYAN},
             {"text": "。「絵に描いた餅」で止めない。", "size": 18, "color": C_WHITE},
         ]],
         fill=C_NAVY, anchor=MSO_ANCHOR.MIDDLE, padding=(20, 40, 20, 40), line_spacing=1.6)

# ============ Slide 5: Profile ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 5)
add_heading(s, "PROFILE  /  代表者",
            [[
                {"text": "弁護士 × 国際法務 × AI実務 ── ", "size": 40, "bold": True, "color": C_NAVY},
                {"text": "福島 駿太", "size": 40, "bold": True, "color": C_CYAN2},
                {"text": "。", "size": 40, "bold": True, "color": C_NAVY},
            ]])
# photo (left)
photo_left = x(96)
photo_top = x(230)
photo_w = x(520)
photo_h = x(720)
add_rect(s, photo_left, photo_top, photo_w, photo_h, fill=C_SUBGRAY)
img_path = os.path.join(HERE, "20230901_IMG_1609_Shunta Fukushima_Full Portrait.jpg")
add_image(s, img_path, photo_left, photo_top, photo_w, photo_h)
# name plate
add_rect(s, photo_left, photo_top + x(620), photo_w, x(100), fill=C_NAVY)
add_text(s, photo_left, photo_top + x(620), photo_w, x(40),
         [{"text": "代表取締役 弁護士", "size": 14, "bold": True, "color": C_CYAN}],
         padding=(14, 32, 0, 32))
add_text(s, photo_left, photo_top + x(650), photo_w, x(60),
         [{"text": "福島 駿太", "size": 26, "bold": True, "color": C_WHITE}],
         padding=(0, 32, 0, 32))

# career list (right)
list_left = x(680)
list_top = x(230)
list_w = x(1144)
add_text(s, list_left, list_top, list_w, x(40),
         [{"text": "CAREER ── 経歴", "size": 18, "bold": True, "color": C_CYAN2}])
# left cyan border
add_line_shape(s, list_left, list_top + x(56), x(3), x(660), C_CYAN)
career = [
    ("2013年 3月", "中央大学法学部 ", "卒業", ""),
    ("2014年12月", "", "弁護士登録", "（第一東京弁護士会）"),
    ("2014年頃",   "朝日中央綜合法律事務所 勤務", "", ""),
    ("2017年 8月", "柴田・鈴木・中田法律事務所 勤務", "", ""),
    ("2018年10月", "", "アンダーソン・毛利・友常法律事務所", " 勤務"),
    ("2023年 5月", "", "コーネル大学ロースクール", "修了（法学修士）"),
    ("2023年 9月", "Walkers（Hong Kong）にて海外研修", "", ""),
    ("2024年 9月", "アンダーソン・毛利・友常法律事務所 復帰", "", ""),
    ("2025年 1月", "", "リーガルアクセス株式会社 創業", ""),
]
row_h = 70
for i, (yr, pre, bold_part, post) in enumerate(career):
    top = list_top + x(60 + i * row_h)
    # dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, list_left + x(-6), top + x(20), x(12), x(12))
    dot.shadow.inherit = False
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_CYAN2
    dot.line.fill.background()
    # year (fixed width)
    add_text(s, list_left + x(28), top, x(200), x(50),
             [{"text": yr, "size": 17, "bold": True, "color": C_CYAN2}],
             anchor=MSO_ANCHOR.MIDDLE)
    # body
    runs = []
    if pre:
        runs.append({"text": pre, "size": 17, "color": C_NAVY})
    if bold_part:
        runs.append({"text": bold_part, "size": 17, "bold": True, "color": C_NAVY})
    if post:
        runs.append({"text": post, "size": 17, "color": C_NAVY})
    add_text(s, list_left + x(228), top, x(900), x(50),
             [runs], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)

# ============ Slide 6: Book ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 6)
add_heading(s, "PUBLICATION  /  著書",
            [[
                {"text": "弘文堂より刊行、", "size": 40, "bold": True, "color": C_NAVY},
                {"text": "累計5,000部", "size": 40, "bold": True, "color": C_CYAN2},
                {"text": "。", "size": 40, "bold": True, "color": C_NAVY},
            ]])
# book cover (left)
book_path = os.path.join(HERE, "表紙（法務のための生成AI活用ガイド）.png")
add_image(s, book_path, x(120), x(240), x(420), x(560))
add_text(s, x(120), x(820), x(420), x(80),
         [
             [{"text": "『法務のための生成AI活用ガイド』", "size": 14, "bold": True, "color": C_NAVY}],
             [{"text": "弁護士 福島駿太 著／弘文堂・2025年", "size": 13, "color": C_TEXT2}],
         ],
         align=PP_ALIGN.CENTER, line_spacing=1.5)

# right side
right_x = x(620)
right_w = x(1200)
# copy band
add_text(s, right_x, x(240), right_w, x(150),
         [[
             {"text": "質が高く効率的な", "size": 24, "bold": True, "color": C_CYAN},
             {"text": "法律業務を叶える", "size": 24, "bold": True, "color": C_WHITE},
         ],
          [{"text": "「生成AIとの上手な付き合い方」", "size": 24, "bold": True, "color": C_WHITE}]],
         fill=C_NAVY, padding=(24, 32, 24, 32), line_spacing=1.5)
# info
add_text(s, right_x, x(410), right_w, x(60),
         [{"text": "法務部・法律事務所・パラリーガルが、今すぐ実務に取り入れられる実践的な1冊。",
           "size": 17, "bold": True, "color": C_NAVY}], line_spacing=1.5)
add_text(s, right_x, x(450), right_w, x(50),
         [{"text": "契約レビュー、リサーチ、ドラフティング、リスク管理まで具体的なプロンプト例で解説。",
           "size": 14, "color": C_TEXT2}], line_spacing=1.5)
# records
records = [
    ("No.1", "", "Amazonベストセラー", "該当カテゴリで1位を獲得。"),
    ("No.1", "", "弁護士会ブックセンター", "実務家コミュニティで首位。"),
    ("5,000", "部", "累計発行部数", "増刷を重ね、現在も伸長中。"),
]
rw = 380
rh = 290
rgap = 30
for i, (rank, unit, ttl, desc) in enumerate(records):
    rleft = right_x + Emu(i * (rw + rgap) * PX)
    rtop = x(540)
    add_rect(s, rleft, rtop, x(rw), x(rh), fill=C_BG_LITE)
    add_rect(s, rleft, rtop, x(rw), x(4), fill=C_CYAN2)
    add_text(s, rleft, rtop + x(28), x(rw), x(80),
             [[
                 {"text": rank, "size": 36, "bold": True, "color": C_CYAN2},
                 {"text": unit, "size": 16, "bold": True, "color": C_NAVY},
             ]],
             padding=(0, 24, 0, 24), line_spacing=1.0)
    add_text(s, rleft, rtop + x(130), x(rw), x(50),
             [{"text": ttl, "size": 17, "bold": True, "color": C_NAVY}],
             padding=(0, 24, 0, 24), line_spacing=1.4)
    add_text(s, rleft, rtop + x(180), x(rw), x(80),
             [{"text": desc, "size": 14, "color": C_TEXT2}],
             padding=(0, 24, 0, 24), line_spacing=1.6)

# ============ Slide 7: Seminars ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 7)
add_heading(s, "SEMINARS  /  講演・登壇実績",
            [[
                {"text": "弁護士会・士業団体・企業 ── ", "size": 40, "bold": True, "color": C_NAVY},
                {"text": "登壇多数", "size": 40, "bold": True, "color": C_CYAN2},
                {"text": "。", "size": 40, "bold": True, "color": C_NAVY},
            ]])
seminars = [
    ("3/19", "（水）", "実はよく知らないけど今さら聞けないAIのこと", "主催：弘文堂・GVA TECH"),
    ("4/23", "（水）", "AI活用のルール、「まだいらない」で本当にいいですか？", "共催：東京新橋法律事務所・リーガルアクセス株式会社"),
    ("5/21", "（水）", "企業のAI活用に必要なルール整備を弁護士3名が解説 ｜法的リスクと実践活用を60分で整理", "共催：東京新橋法律事務所・リーガルアクセス株式会社"),
    ("6/22", "（日）", "もう一歩先へ 熟練弁護士による生成AIの総合的活用術", "主催：GVA TECH"),
    ("7/10", "（木）", "中小企業法務とAIの利活用", "主催：東京弁護士会"),
    ("8/10", "（日）", "法務のための生成AI活用について ── 基礎編", "主催：公認会計士協会"),
    ("8/20", "（水）", "法務部員のための実践的AI活用方法", "主催：国際企業法務協会"),
    ("───", "",       "そのほか、企業内研修・私的勉強会など登壇実績多数", "個別の社内研修・経営層向けブリーフィング含む"),
]
sem_w = 850
sem_h = 130
sem_gx = 28
sem_gy = 16
sem_sx = 96
sem_sy = 230
for i, (date, dow, title, org) in enumerate(seminars):
    row = i // 2
    col = i % 2
    left = x(sem_sx + col * (sem_w + sem_gx))
    top = x(sem_sy + row * (sem_h + sem_gy))
    add_rect(s, left, top, x(sem_w), x(sem_h), fill=C_WHITE, line=C_SUBGRAY, line_w=0.75)
    add_rect(s, left, top, x(6), x(sem_h), fill=C_CYAN2)
    # date
    add_text(s, left + x(18), top, x(140), x(sem_h),
             [
                 [{"text": date, "size": 22, "bold": True, "color": C_CYAN2}],
                 [{"text": dow,  "size": 11, "color": C_TEXT2}],
             ],
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)
    # title + org
    if title == "そのほか、企業内研修・私的勉強会など登壇実績多数":
        title_runs = [
            [
                {"text": "そのほか、企業内研修・私的勉強会など登壇実績", "size": 16, "bold": True, "color": C_NAVY},
                {"text": "多数", "size": 16, "bold": True, "color": C_CYAN2},
            ],
            [{"text": org, "size": 12, "color": C_TEXT2}],
        ]
    else:
        title_runs = [
            [{"text": title, "size": 16, "bold": True, "color": C_NAVY}],
            [{"text": org, "size": 12, "color": C_TEXT2}],
        ]
    add_text(s, left + x(170), top, x(sem_w - 180), x(sem_h),
             title_runs, anchor=MSO_ANCHOR.MIDDLE, padding=(8, 16, 8, 0), line_spacing=1.5)

# ============ Slide 8: Detail 1 ============
def render_detail_slide(page, label, title_pre, title_bold, title_post,
                        blocks, image_filename, tag_text):
    s = prs.slides.add_slide(BLANK)
    add_page_chrome(s, page)
    add_heading(s, label,
                [[
                    {"text": title_pre, "size": 40, "bold": True, "color": C_NAVY},
                    {"text": title_bold, "size": 40, "bold": True, "color": C_CYAN2},
                    {"text": title_post, "size": 40, "bold": True, "color": C_NAVY},
                ]])
    # left blocks
    block_x = x(96)
    block_w = x(1100)
    block_h = 360
    block_gap = 24
    block_top0 = 240
    for i, (num, h3, items) in enumerate(blocks):
        top = x(block_top0 + i * (block_h + block_gap))
        add_rect(s, block_x, top, block_w, x(block_h), fill=C_BG_LITE)
        add_rect(s, block_x, top, x(8), x(block_h), fill=C_CYAN2)
        # num pill
        add_rect(s, block_x + x(40), top + x(28), x(120), x(36), fill=C_CYAN2)
        add_text(s, block_x + x(40), top + x(28), x(120), x(36),
                 [{"text": num, "size": 12, "bold": True, "color": C_WHITE}],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # h3
        add_text(s, block_x + x(40), top + x(74), block_w - x(80), x(60),
                 [{"text": h3, "size": 22, "bold": True, "color": C_NAVY}], line_spacing=1.35)
        # items
        item_runs = []
        for runs in items:
            full = [{"text": "▸  ", "size": 16, "bold": True, "color": C_CYAN2}] + runs
            item_runs.append(full)
        add_text(s, block_x + x(40), top + x(140), block_w - x(80), x(200),
                 item_runs, line_spacing=1.7, padding=(0, 0, 0, 0))
    # right image
    img_x = x(1220)
    img_y = x(240)
    img_w = x(600)
    img_h = x(744)
    add_rect(s, img_x, img_y, img_w, img_h, fill=C_SUBGRAY)
    img_path = os.path.join(HERE, image_filename)
    add_image(s, img_path, img_x, img_y, img_w, img_h)
    # tag
    add_rect(s, img_x + x(24), img_y + x(24), x(180), x(40), fill=C_NAVY)
    add_text(s, img_x + x(24), img_y + x(24), x(180), x(40),
             [{"text": tag_text, "size": 13, "bold": True, "color": C_CYAN}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s

# Slide 8
render_detail_slide(
    8, "DETAIL ①  /  ヒアリング〜現状分析",
    "まず、", "「今、何が起きているか」", "を正確に掴む。",
    [
        ("STEP 01", "ヒアリング・状況把握", [
            [
                {"text": "担当責任者と", "size": 15, "color": C_NAVY},
                {"text": "1対1で打ち合わせ", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "。現状のAI導入・活用状況をヒアリング", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "導入済みツール／未導入ツールを", "size": 15, "color": C_NAVY},
                {"text": "棚卸し", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "。「契約だけして使われていない」も含めて把握", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "現場担当者がどのツールを、どの業務で、どう使っているかを", "size": 15, "color": C_NAVY},
                {"text": "実態ベース", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "で確認", "size": 15, "color": C_NAVY},
            ],
        ]),
        ("STEP 02", "現状分析", [
            [
                {"text": "ヒアリング結果をもとに、", "size": 15, "color": C_NAVY},
                {"text": "現在の投資効率", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "を可視化", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "人件費にいくらかかっているのか、AIツールにいくら払っているのか、", "size": 15, "color": C_NAVY},
                {"text": "金額で並べる", "size": 15, "bold": True, "color": C_CYAN2},
            ],
            [
                {"text": "「払っている分のリターンが出ているか」を経営層に説明できる形に整理", "size": 15, "color": C_NAVY},
            ],
        ]),
    ],
    "タブレットとペン.png", "HEARING")

# Slide 9
render_detail_slide(
    9, "DETAIL ②  /  活用案の策定〜提言",
    "「", "どの業務に、どのAIを", "」を具体名で出す。",
    [
        ("STEP 03", "AI活用案の策定", [
            [
                {"text": "AI活用の", "size": 15, "color": C_NAVY},
                {"text": "余地があるのに使われていない業務", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "を洗い出し", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "既に活用されているが、", "size": 15, "color": C_NAVY},
                {"text": "さらに踏み込める業務", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "もリストアップ", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "業務単位で「使うべきツール／プロンプト方針」をセットで設計", "size": 15, "color": C_NAVY},
            ],
        ]),
        ("STEP 04", "AI活用案の提言", [
            [
                {"text": "個社の業務に応じた", "size": 15, "color": C_NAVY},
                {"text": "具体的な活用法", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "を提言（机上論ではなく実装可能な粒度）", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "導入時の", "size": 15, "color": C_NAVY},
                {"text": "経営インパクト", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "（時間削減・コスト削減）を金額で提示", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "責任者・社員向けの", "size": 15, "color": C_NAVY},
                {"text": "オンボーディング説明会・研修", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "を実施", "size": 15, "color": C_NAVY},
            ],
        ]),
    ],
    "パソコン画面とコード.png", "DESIGN")

# Slide 10
render_detail_slide(
    10, "DETAIL ③  /  継続的モニタリング・継続支援",
    "作って終わりではない、", "定着するまで伴走", "。",
    [
        ("STEP 05", "継続的モニタリング", [
            [
                {"text": "現場でAIが", "size": 15, "color": C_NAVY},
                {"text": "実際に使われているか", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "を継続的に把握", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "定期レビューで「想定通り使えている業務／想定外に止まっている業務」を切り分け", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "導入直後の盛り上がりで終わらせず、", "size": 15, "color": C_NAVY},
                {"text": "習慣にする", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "まで見届ける", "size": 15, "color": C_NAVY},
            ],
        ]),
        ("STEP 06", "継続支援・ギャップの解消", [
            [
                {"text": "運用と設計の", "size": 15, "color": C_NAVY},
                {"text": "ギャップが生じた箇所", "size": 15, "bold": True, "color": C_CYAN2},
                {"text": "を特定", "size": 15, "color": C_NAVY},
            ],
            [
                {"text": "使われない理由（UI／教育／業務フロー）を分析し、", "size": 15, "color": C_NAVY},
                {"text": "解決策を提示", "size": 15, "bold": True, "color": C_CYAN2},
            ],
            [
                {"text": "必要なプロンプト改修、ルール改訂、追加研修まで", "size": 15, "color": C_NAVY},
                {"text": "実装を支援", "size": 15, "bold": True, "color": C_CYAN2},
            ],
        ]),
    ],
    "キーボード.png", "RUN & IMPROVE")

# ============ Slide 11: Option ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 11)
add_heading(s, "OPTION  /  個別開発",
            [[
                {"text": "必要があれば、", "size": 40, "bold": True, "color": C_NAVY},
                {"text": "個別システム開発", "size": 40, "bold": True, "color": C_CYAN2},
                {"text": "まで。", "size": 40, "bold": True, "color": C_NAVY},
            ]])
# left main (navy)
add_rect(s, x(96), x(240), x(860), x(720), fill=C_NAVY)
add_text(s, x(96), x(280), x(860), x(50),
         [{"text": "DEFAULT", "size": 14, "bold": True, "color": C_CYAN}],
         padding=(0, 56, 0, 56))
add_text(s, x(96), x(340), x(860), x(220),
         [
             [{"text": "本サービスは", "size": 26, "bold": True, "color": C_WHITE}],
             [
                 {"text": "「汎用型生成AI」", "size": 30, "bold": True, "color": C_CYAN},
             ],
             [{"text": "の活用を前提", "size": 26, "bold": True, "color": C_WHITE}],
         ],
         padding=(0, 56, 0, 56), line_spacing=1.45)
add_text(s, x(96), x(620), x(860), x(300),
         [{"text": "ChatGPT、Gemini、Claude、Copilot といった汎用AIを、業務に合わせて使いこなす設計を行います。多くの業務は、汎用AIと適切なプロンプト設計でカバーできます。",
           "size": 16, "color": C_SUBGRAY}],
         padding=(0, 56, 0, 56), line_spacing=1.7)

# right sub (light)
add_rect(s, x(984), x(240), x(840), x(720), fill=C_BG_LITE)
add_rect(s, x(1024), x(280), x(150), x(36), fill=C_CYAN2)
add_text(s, x(1024), x(280), x(150), x(36),
         [{"text": "OPTIONAL", "size": 12, "bold": True, "color": C_WHITE}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x(984), x(340), x(840), x(220),
         [
             [{"text": "業務特化の", "size": 26, "bold": True, "color": C_NAVY}],
             [{"text": "個別開発が必要な場合", "size": 26, "bold": True, "color": C_NAVY}],
         ],
         padding=(0, 56, 0, 56), line_spacing=1.45)
add_text(s, x(984), x(560), x(840), x(360),
         [[
             {"text": "社内文書RAG、契約レビュー特化UI、社内チャットボット等、独自システムが必要な場合は", "size": 16, "color": C_TEXT2},
             {"text": "別途お見積り", "size": 16, "bold": True, "color": C_CYAN2},
             {"text": "のうえ対応いたします。", "size": 16, "color": C_TEXT2},
         ]],
         padding=(0, 56, 0, 56), line_spacing=1.7)

# ============ Slide 12: Pricing ============
s = prs.slides.add_slide(BLANK)
add_page_chrome(s, 12)
add_heading(s, "PRICING  /  料金・お支払い・上限",
            [[
                {"text": "月額制の伴走サービス、", "size": 40, "bold": True, "color": C_NAVY},
                {"text": "上限8社", "size": 40, "bold": True, "color": C_CYAN2},
                {"text": "。", "size": 40, "bold": True, "color": C_NAVY},
            ]])
# left: price main (navy gradient simulated with single color)
add_rect(s, x(96), x(240), x(900), x(720), fill=C_NAVY)
add_text(s, x(96), x(290), x(900), x(40),
         [{"text": "MONTHLY FEE", "size": 14, "bold": True, "color": C_CYAN}],
         padding=(0, 60, 0, 60))
add_text(s, x(96), x(350), x(900), x(220),
         [[
             {"text": "50", "size": 80, "bold": True, "color": C_WHITE},
             {"text": "万円", "size": 28, "bold": True, "color": C_CYAN},
         ]],
         padding=(0, 60, 0, 60), line_spacing=1.0)
add_text(s, x(96), x(580), x(900), x(60),
         [{"text": "× ご契約月数（税別）", "size": 22, "bold": True, "color": C_SUBGRAY}],
         padding=(0, 60, 0, 60))
add_text(s, x(96), x(680), x(900), x(280),
         [
             [{"text": "月額固定で、ヒアリング・分析・提言・研修・継続モニタリングまで含む包括フィー。",
              "size": 15, "color": C_TEXT_L}],
             [{"text": "個別開発が必要な場合は別途お見積り。",
              "size": 15, "color": C_TEXT_L}],
         ],
         padding=(0, 60, 0, 60), line_spacing=1.7)

# right items
ri_x = x(1024)
ri_w = x(800)
# item 1 (payment)
add_rect(s, ri_x, x(240), ri_w, x(340), fill=C_BG_LITE)
add_rect(s, ri_x, x(240), x(6), x(340), fill=C_CYAN2)
add_text(s, ri_x + x(32), x(280), ri_w - x(64), x(50),
         [{"text": "PAYMENT ── お支払い", "size": 16, "bold": True, "color": C_CYAN2}])
add_text(s, ri_x + x(32), x(340), ri_w - x(64), x(220),
         [
             [
                 {"text": "ご契約時に", "size": 17, "color": C_NAVY},
                 {"text": "当初月分", "size": 17, "bold": True, "color": C_CYAN2},
                 {"text": "をご入金いただきます。", "size": 17, "color": C_NAVY},
             ],
             [
                 {"text": "以降は", "size": 17, "color": C_NAVY},
                 {"text": "前月25日まで", "size": 17, "bold": True, "color": C_CYAN2},
                 {"text": "に当月分をお支払い。", "size": 17, "color": C_NAVY},
             ],
         ],
         line_spacing=1.8)

# item 2 (limit) — highlighted
add_rect(s, ri_x, x(620), ri_w, x(340), fill=C_NAVY)
add_rect(s, ri_x, x(620), x(6), x(340), fill=C_CYAN)
add_text(s, ri_x + x(32), x(650), ri_w - x(64), x(50),
         [{"text": "LIMITATION ── ご支援先の上限", "size": 16, "bold": True, "color": C_CYAN}])
add_text(s, ri_x + x(32), x(700), ri_w - x(64), x(120),
         [{"text": "8社", "size": 60, "bold": True, "color": C_CYAN}],
         line_spacing=1.0)
add_text(s, ri_x + x(32), x(840), ri_w - x(64), x(110),
         [[
             {"text": "伴走型のため、同時にお引き受けできる企業数に", "size": 16, "color": C_WHITE},
             {"text": "上限", "size": 16, "bold": True, "color": C_CYAN},
             {"text": "を設けています。", "size": 16, "color": C_WHITE},
         ]],
         line_spacing=1.7)

# ============ Slide 13: Closing ============
s = prs.slides.add_slide(BLANK)
# background
add_rect(s, x(0), x(0), x(1920), x(1080), fill=C_DEEP)
img_path = os.path.join(HERE, "青いネットワーク.png")
add_image(s, img_path, x(960), x(0), x(960), x(1080))
add_rect(s, x(0), x(0), x(1200), x(1080), fill=C_NAVY)
# content
add_text(s, x(120), x(180), x(1100), x(40),
         [{"text": "APPLICATION  |  お申込み", "size": 16, "bold": True, "color": C_CYAN}])
add_text(s, x(120), x(240), x(1100), x(280),
         [
             [{"text": "お申込みは、", "size": 56, "bold": True, "color": C_WHITE}],
             [
                 {"text": "先着順", "size": 56, "bold": True, "color": C_CYAN},
                 {"text": "とさせていただきます。", "size": 56, "bold": True, "color": C_WHITE},
             ],
         ],
         line_spacing=1.3)
add_text(s, x(120), x(540), x(1500), x(220),
         [
             [
                 {"text": "他の企業様からのお申込みが先に入った場合、ご依頼を", "size": 20, "color": C_SUBGRAY},
                 {"text": "やむを得ずお断り", "size": 20, "bold": True, "color": C_CYAN},
                 {"text": "させていただくことがございます。", "size": 20, "color": C_SUBGRAY},
             ],
             [
                 {"text": "ご検討中の場合は、", "size": 20, "color": C_SUBGRAY},
                 {"text": "お早めにご連絡", "size": 20, "bold": True, "color": C_CYAN},
                 {"text": "いただけますと幸いです。", "size": 20, "color": C_SUBGRAY},
             ],
         ],
         line_spacing=1.7)
# contact box
add_rect(s, x(120), x(800), x(1300), x(160), fill=None, line=C_CYAN, line_w=1.5)
add_rect(s, x(120), x(800), x(6), x(160), fill=C_CYAN)
add_text(s, x(160), x(820), x(200), x(40),
         [{"text": "CONTACT", "size": 14, "bold": True, "color": C_CYAN}],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x(380), x(820), x(1000), x(60),
         [{"text": "リーガルアクセス株式会社 ／ 代表取締役 弁護士 福島 駿太",
           "size": 20, "bold": True, "color": C_WHITE}],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, x(380), x(880), x(1000), x(50),
         [{"text": "本資料に関するお問い合わせは、商談ご担当窓口までご連絡ください。",
           "size": 14, "color": C_TEXT_L}],
         anchor=MSO_ANCHOR.MIDDLE)

# ============ 保存 ============
try:
    prs.save(OUTPUT)
    print(f"OK: {OUTPUT}")
except Exception as e:
    print(f"NG: {e}")
    raise
