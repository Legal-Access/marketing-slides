# -*- coding: utf-8 -*-
"""
社内勉強会スライド（プレゼン版・16:9）

- 目的: 既存の「特別レポート」由来の読み物スライド（本文ベタ書き24枚）を、
        話者が見せながら喋るプレゼンスライド（キーワード・カード・写真主役）に作り替える。
- 元データ: 20260529_保険代理店向け_社内勉強会スライド_補足資料.pptx
            （本文テキストの原稿・埋め込み写真8枚の抽出元）
- 出力: 同名ファイルへ上書き（元はGit履歴に残る）
- 方針: 事実は一切創作しない。下記の文言はすべて元pptx本文からの抽出・短縮であり新事実は足さない。
        最小フォント縛り（紙向けの16pt）は撤廃。見出し44〜48 / 本文24〜28 / ラベル18〜20。
"""
import io
import os
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 定数

BASE = r"g:\共有ドライブ\マーケティング部\23_保険代理店開拓"
SRC  = os.path.join(BASE, "20260529_保険代理店向け_社内勉強会スライド_補足資料.pptx")
OUT  = SRC  # 元ファイル名で上書き

# 16:9
SW = Inches(13.333)
SH = Inches(7.5)

# 配色（元レポートを継承）
GREEN      = RGBColor(0x2B, 0x64, 0x44)
GREEN_DARK = RGBColor(0x1E, 0x4D, 0x33)
GREEN_LT   = RGBColor(0xE8, 0xF2, 0xEB)
GREEN_MID  = RGBColor(0xC8, 0xDF, 0xD0)
TEXT       = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SUB   = RGBColor(0x55, 0x55, 0x55)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xB8, 0x96, 0x4A)

SANS  = "Noto Sans JP"
SERIF = "Noto Serif JP"

# プレゼン用フォントサイズ（紙の最小16pt縛りを撤廃）
H1   = 48   # 大見出し
H2   = 40   # セクション見出し
H3   = 28   # 小見出し
BODY = 26   # 本文（短文）
LEAD = 22   # リード
LBL  = 18   # ラベル

# ---------------------------------------------------------------- 画像抽出（_html_imgが消えているのでpptx内から再利用）

def _load_images_from_pptx(path):
    """元pptxの ppt/media/imageN.jpg を bytes で読み出し、px寸法とともに保持。
       slide→media対応はサイズ・配置順から確定済み（下のキーへ手当て）。"""
    z = zipfile.ZipFile(path)
    blobs = {}
    for n in z.namelist():
        if n.startswith("ppt/media/") and n.lower().endswith((".jpg", ".jpeg", ".png")):
            blobs[os.path.basename(n)] = z.read(n)
    z.close()
    # 確定対応（調査で確認）：
    #   image1=cover / image2=依頼者 / image3=書類対応 / image4=法律書籍
    #   image5=評価損 / image6=家族 / image7=相談 / image8=cta
    mapping = {
        "cover":   "image1.jpg",
        "依頼者":  "image2.jpg",
        "書類対応": "image3.jpg",
        "法律書籍": "image4.jpg",
        "評価損":  "image5.jpg",
        "家族":    "image6.jpg",
        "相談":    "image7.jpg",
        "cta":     "image8.jpg",
    }
    out = {}
    for alt, fname in mapping.items():
        b = blobs.get(fname)
        if b is None:
            continue
        w, h = _img_px_bytes(b)
        out[alt] = (b, w, h)
    return out

def _img_px_bytes(data):
    """JPEG/PNGのpx寸法を外部ライブラリなしでbytesから読む。"""
    import struct
    try:
        if data[:8] == b"\x89PNG\r\n\x1a\n":
            return struct.unpack(">II", data[16:24])
        if data[:2] == b"\xff\xd8":
            i = 2; n = len(data)
            while i < n:
                if data[i] != 0xFF:
                    i += 1; continue
                mk = data[i + 1]
                if mk in (0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7, 0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF):
                    h = struct.unpack(">H", data[i + 5:i + 7])[0]
                    w = struct.unpack(">H", data[i + 7:i + 9])[0]
                    return w, h
                i += 2 + struct.unpack(">H", data[i + 2:i + 4])[0]
    except Exception:
        pass
    return None, None

# ---------------------------------------------------------------- 基盤

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

IMGS = _load_images_from_pptx(SRC)

def _set_ea_font(run, name):
    """日本語（East Asian）フォントを確実に効かせる。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.25, space_after=Pt(6), wrap=True):
    """runs: list of paragraphs; 各paragraph = list of (text, size, color, bold, font)。"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_after = space_after
        p.space_before = Pt(0)
        for (t, size, color, bold, font) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            _set_ea_font(r, font)
    return tb

def P(text, size=BODY, color=TEXT, bold=False, font=SANS):
    """1ラン1段落の糖衣。"""
    return [(text, size, color, bold, font)]

def add_picture_cover(s, alt, x, y, w, h):
    """枠(w,h)をアスペクト比維持で覆い、はみ出しはクロップ。bytesから配置。"""
    blob, iw, ih = IMGS[alt]
    stream = io.BytesIO(blob)
    if not iw:
        return s.shapes.add_picture(stream, x, y, width=w, height=h)
    fw, fh = float(w), float(h)
    fr = fw / fh
    ir = iw / ih
    if ir > fr:
        new_h = fh; new_w = fh * ir
    else:
        new_w = fw; new_h = fw / ir
    px = x - (new_w - fw) / 2
    py = y - (new_h - fh) / 2
    pic = s.shapes.add_picture(stream, int(px), int(py), int(new_w), int(new_h))
    pic.crop_left = max(0, (new_w - fw) / 2 / new_w)
    pic.crop_right = max(0, (new_w - fw) / 2 / new_w)
    pic.crop_top = max(0, (new_h - fh) / 2 / new_h)
    pic.crop_bottom = max(0, (new_h - fh) / 2 / new_h)
    pic.left = int(x); pic.top = int(y); pic.width = int(w); pic.height = int(h)
    return pic

def section_label(s, x, y, text, color=GREEN):
    """英字の小見出しラベル。"""
    txt(s, x, y, Inches(8), Inches(0.5),
        [P(text.upper(), LBL, color, True, SANS)], line_spacing=1.0, space_after=Pt(0))

def heading(s, x, y, w, text, size=H2, color=GREEN_DARK):
    return txt(s, x, y, w, Inches(1.2),
               [P(text, size, color, True, SERIF)], line_spacing=1.1, space_after=Pt(0))

def h3(s, x, y, w, text, size=H3, color=GREEN_DARK):
    return txt(s, x, y, w, Inches(0.8),
               [P(text, size, color, True, SERIF)], line_spacing=1.15, space_after=Pt(0))

def num_badge(s, x, y, n, dia=Inches(0.7), fill=GREEN, fg=WHITE, fs=30):
    sp = rect(s, x, y, dia, dia, fill=fill, shape=MSO_SHAPE.OVAL)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = fg
    _set_ea_font(r, SERIF)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sp

def case_band(s, case_no, title_lines, lead):
    """CASE扉の濃緑帯（上部）＋大番号＋タイトル＋リード一行。"""
    rect(s, 0, 0, SW, Inches(2.4), fill=GREEN_DARK)
    rect(s, Inches(10.6), 0, Inches(2.73), Inches(2.4), fill=GREEN)
    txt(s, Inches(0.7), Inches(0.45), Inches(5), Inches(0.6),
        [P(case_no, 24, GREEN_MID, True, SANS)], line_spacing=1.0, space_after=Pt(0))
    txt(s, Inches(0.7), Inches(1.0), Inches(12.0), Inches(1.3),
        [P(l, 34, WHITE, True, SERIF) for l in title_lines],
        line_spacing=1.12, space_after=Pt(0))
    # リード（帯の下・1行）
    txt(s, Inches(0.7), Inches(3.0), Inches(12.0), Inches(1.2),
        lead, line_spacing=1.4, space_after=Pt(0))

# ---- 追加ヘルパ（プレゼン用パーツ） ----

def kw_card(s, x, y, w, h, head_runs, body=None, fill=GREEN_LT, top_accent=GREEN,
            head_size=H3, body_size=20, body_color=TEXT):
    """キーワードカード：上に強調見出し、下に一言補足。"""
    rect(s, x, y, w, h, fill=fill)
    rect(s, x, y, w, Inches(0.12), fill=top_accent)  # 上アクセント
    pad = Inches(0.3)
    inner_x = Emu(int(x) + int(pad))
    inner_w = Emu(int(w) - 2 * int(pad))
    txt(s, inner_x, Emu(int(y) + int(Inches(0.3))), inner_w, Inches(0.9),
        [head_runs], line_spacing=1.15, space_after=Pt(0))
    if body:
        txt(s, inner_x, Emu(int(y) + int(Inches(1.15))), inner_w, Emu(int(h) - int(Inches(1.4))),
            [P(body, body_size, body_color, False, SANS)], line_spacing=1.3, space_after=Pt(0))

def flow_step(s, x, y, w, h, label_runs, fill=GREEN_LT, txt_color=GREEN_DARK):
    """横フローの1ステップ（角丸ボックス）。"""
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Emu(int(x) + int(Inches(0.15))), y, Emu(int(w) - 2 * int(Inches(0.15))), h,
        [label_runs], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.2, space_after=Pt(0))

def arrow(s, x, y, w, h, fill=GREEN):
    rect(s, x, y, w, h, fill=fill, shape=MSO_SHAPE.RIGHT_ARROW)

def big_banner(s, x, y, w, h, runs, fill=GREEN, anchor=MSO_ANCHOR.MIDDLE):
    """全面・帯メッセージ。"""
    rect(s, x, y, w, h, fill=fill)
    txt(s, Emu(int(x) + int(Inches(0.5))), y, Emu(int(w) - 2 * int(Inches(0.5))), h,
        runs, align=PP_ALIGN.CENTER, anchor=anchor, line_spacing=1.3, space_after=Pt(0))

# 強調runの糖衣：重要語だけGREEN太字、地の文は通常
def em(strong, normal_before="", normal_after="", size=BODY, strong_color=GREEN_DARK,
       normal_color=TEXT, font_strong=SANS, font_normal=SANS):
    para = []
    if normal_before:
        para.append((normal_before, size, normal_color, False, font_normal))
    para.append((strong, size, strong_color, True, font_strong))
    if normal_after:
        para.append((normal_after, size, normal_color, False, font_normal))
    return para

# ================================================================ スライド構築

def s01_cover():
    s = slide()
    rect(s, 0, 0, Inches(7.8), SH, fill=WHITE)
    add_picture_cover(s, "cover", Inches(7.8), 0, Inches(5.533), SH)
    section_label(s, Inches(0.8), Inches(0.8), "保険代理店向け 社内勉強会")
    # タグ
    rect(s, Inches(0.8), Inches(1.45), Inches(3.6), Inches(0.6), fill=GREEN_LT)
    txt(s, Inches(0.8), Inches(1.45), Inches(3.6), Inches(0.6),
        [P("交通事故 × 弁護士費用特約", LBL, GREEN_DARK, True, SANS)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
    # 大タイトル
    txt(s, Inches(0.8), Inches(2.5), Inches(6.8), Inches(2.4),
        [P("弁護士費用特約は", H1, GREEN_DARK, True, SERIF),
         [("「揉める前」", H1, GREEN, True, SERIF), ("に動く", H1, GREEN_DARK, True, SERIF)]],
        line_spacing=1.18, space_after=Pt(0))
    txt(s, Inches(0.8), Inches(5.1), Inches(6.8), Inches(0.8),
        [P("3つの実例で考える、初動から弁護士を入れる意義", LEAD, TEXT_SUB, False, SANS)],
        line_spacing=1.3, space_after=Pt(0))
    # 署名
    txt(s, Inches(0.8), Inches(6.5), Inches(6.8), Inches(0.7),
        [[("弁護士法人 相生綜合法律事務所　", LBL, GREEN_DARK, True, SERIF),
          ("代表弁護士 福島駿太", LBL, TEXT_SUB, False, SANS)]],
        line_spacing=1.1, space_after=Pt(0))

def s02_intro():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Introduction")
    heading(s, Inches(0.7), Inches(1.1), Inches(8), "弁護士には、いつ相談すべきか")
    add_picture_cover(s, "依頼者", Inches(9.3), Inches(0.6), Inches(3.4), Inches(3.0))
    # 問い3カード（縦）
    qs = [
        em("揉めてから", normal_after="で十分？"),
        em("揉めていなくても", normal_after="頼んでいい？"),
        [("特約", H3, GREEN, True, SANS), ("が使えるなら？", H3, GREEN_DARK, True, SANS)],
    ]
    y = Inches(2.7)
    for q in qs:
        rect(s, Inches(0.7), y, Inches(8.0), Inches(0.95), fill=GREEN_LT)
        txt(s, Inches(1.0), y, Inches(7.5), Inches(0.95), [q],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.1)))
    # 答え（強調帯）
    big_banner(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.95),
               [[("答えは——", 24, WHITE, False, SANS),
                 ("できるだけ早く", 30, WHITE, True, SERIF)]])

def s03_conclusion_first():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Today's Point")
    heading(s, Inches(0.7), Inches(1.1), Inches(12), "今日の結論を、先に")
    # 中央バナー
    big_banner(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.3),
               [[("初動から弁護士", 34, WHITE, True, SERIF),
                 ("　が三方よし", 26, WHITE, False, SANS)],
                [("依頼者 ／ 紹介元 ／ 相手方保険会社", 20, GREEN_MID, True, SANS)]])
    # 3CASE目次
    items = [
        ("01", em("治療費打切り", normal_after="後の", size=24) + [("請求漏れ", 24, GREEN, True, SANS)]),
        ("02", [("後遺障害", 24, TEXT, False, SANS), ("診断書の修正", 24, GREEN, True, SANS), ("で認定", 24, TEXT, False, SANS)]),
        ("03", [("外車の", 24, TEXT, False, SANS), ("評価損", 24, GREEN, True, SANS), ("の見落とし", 24, TEXT, False, SANS)]),
    ]
    y = Inches(4.0)
    for num, runs in items:
        num_badge(s, Inches(0.9), y, num, dia=Inches(0.85), fs=28)
        txt(s, Inches(2.1), y, Inches(10.5), Inches(0.85), [runs],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.05)))

# ---- CASE 01 ----
def s04_case01_door():
    s = slide()
    case_band(s, "CASE 01", ["治療費打切り後に請求漏れが見つかった事例"],
              [[("打切りを告げられて初めて相談 → ", LEAD, TEXT_SUB, False, SANS),
                ("本来請求できる項目が相当程度漏れている", LEAD, GREEN_DARK, True, SANS)]])
    # 下部に状況キーワード
    kw_card(s, Inches(0.7), Inches(4.4), Inches(12.0), Inches(2.4),
            [("むち打ち等で通院中、相手方保険会社から", 22, TEXT, False, SANS),
             ("「今月末で治療費を打切ります」", 22, GREEN, True, SANS)],
            body="——この段階で初めて来られるケースが少なくない。内容を確認すると、請求項目が漏れていることがある。",
            body_size=22)

def s05_case01_leak():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Case 01")
    heading(s, Inches(0.7), Inches(1.1), Inches(12), "実際に漏れていたもの")
    # 上バナー
    big_banner(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(0.85),
               [[("やり取りは ", 22, WHITE, False, SANS),
                 ("治療費の話だけ", 26, WHITE, True, SERIF),
                 (" になっていた", 22, WHITE, False, SANS)]], fill=GREEN_DARK)
    # 2×2カード
    cards = [
        ([("入通院慰謝料", H3, GREEN_DARK, True, SANS)], "通院期間に応じた精神的損害の補償"),
        ([("交通費", H3, GREEN_DARK, True, SANS)], "通院時の交通機関利用料・駐車場代等"),
        ([("休業損害", H3, GREEN_DARK, True, SANS)], "事故による収入減少の補償"),
        ([("営業損害", H3, GREEN_DARK, True, SANS)], "自営業者・経営者に特有の損害項目"),
    ]
    cw = Inches(5.85); ch = Inches(1.55); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.7); y0 = Inches(3.35)
    for i, (head, body) in enumerate(cards):
        col = i % 2; row = i // 2
        x = Emu(int(x0) + col * (int(cw) + int(gx)))
        y = Emu(int(y0) + row * (int(ch) + int(gy)))
        kw_card(s, x, y, cw, ch, head, body=body, body_size=18)

def s06_case01_lesson():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Case 01")
    heading(s, Inches(0.7), Inches(1.1), Inches(8), "弁護士が入ると")
    add_picture_cover(s, "書類対応", Inches(9.3), Inches(0.6), Inches(3.4), Inches(3.4))
    steps = [
        em("治療継続", normal_after="の意見書・診断書 → 期間延長を求める", size=22),
        [("漏れた項目を ", 22, TEXT, False, SANS), ("まとめて請求", 22, GREEN, True, SANS)],
        [("結果として ", 22, TEXT, False, SANS), ("適切な解決", 22, GREEN, True, SANS), ("へ", 22, TEXT, False, SANS)],
    ]
    y = Inches(2.4)
    for i, st in enumerate(steps, 1):
        rect(s, Inches(0.7), y, Inches(8.0), Inches(0.9), fill=GREEN_LT)
        num_badge(s, Emu(int(Inches(0.85))), Emu(int(y) + int(Inches(0.15))), i, dia=Inches(0.6), fs=22)
        txt(s, Inches(1.7), y, Inches(6.9), Inches(0.9), [st],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.05)))
    # プルクオート
    big_banner(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(1.25),
               [P("後から直すより、最初から入る方がシンプルで早い。", 28, WHITE, True, SERIF)])

# ---- CASE 02 ----
def s07_case02_door():
    s = slide()
    case_band(s, "CASE 02", ["後遺障害診断書の記載を修正し", "後遺障害認定が得られた事例"],
              [[("診断書は出せば認定されるわけではない —— ", LEAD, TEXT_SUB, False, SANS),
                ("どう書くかが結果を左右する", LEAD, GREEN_DARK, True, SANS)]])
    kw_card(s, Inches(0.7), Inches(4.4), Inches(12.0), Inches(2.4),
            [("保険代理店からのご紹介。相談時点で", 22, TEXT, False, SANS),
             ("診断書は提出済み", 22, GREEN, True, SANS)],
            body="認定機関に提出済みの状況。まずは結果を待つしかなかった。",
            body_size=22)

def s08_case02_flow():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Case 02")
    heading(s, Inches(0.7), Inches(1.1), Inches(8.3), "いったん不認定。そこから逆転")
    add_picture_cover(s, "法律書籍", Inches(9.5), Inches(0.6), Inches(3.2), Inches(3.0))
    # 縦フロー4ステップ（左半分）
    steps = [
        [("提出済みで相談 → ", 22, TEXT, False, SANS), ("不認定", 24, RGBColor(0xB0,0x3A,0x2E), True, SANS)],
        em("弁護士介入", normal_after="・整形外科医と協議", size=22),
        em("記載を適切に修正", normal_after="し再提出", size=22),
        [("→ ", 22, TEXT, False, SANS), ("認定", 28, GREEN, True, SERIF), ("を獲得", 22, TEXT, False, SANS)],
    ]
    y = Inches(2.6)
    for i, st in enumerate(steps):
        flow_step(s, Inches(0.7), y, Inches(8.3), Inches(0.85), st)
        if i < len(steps) - 1:
            arrow(s, Inches(4.55), Emu(int(y) + int(Inches(0.87))), Inches(0.6), Inches(0.28))
        y = Emu(int(y) + int(Inches(1.15)))

def s09_case02_lesson():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Case 02")
    heading(s, Inches(0.7), Inches(1.1), Inches(12), "診断書の記載が、結果を左右する")
    big_banner(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(1.0),
               [[("どの症状を・どう書くか", 30, WHITE, True, SERIF),
                 (" が認定を決める", 24, WHITE, False, SANS)]], fill=GREEN_DARK)
    pts = [
        em("認定基準を踏まえた記載のポイント", normal_after="を医師へ伝える", size=22),
        em("必要な検査・画像資料", normal_after="を初期段階から整備する", size=22),
        em("異議申立", normal_before="不認定時の", normal_after="も視野に並行準備する", size=22),
    ]
    y = Inches(3.5)
    for i, pt in enumerate(pts, 1):
        rect(s, Inches(0.7), y, Inches(12.0), Inches(0.85), fill=GREEN_LT)
        num_badge(s, Emu(int(Inches(0.9))), Emu(int(y) + int(Inches(0.13))), i, dia=Inches(0.58), fs=22)
        txt(s, Inches(1.8), y, Inches(10.7), Inches(0.85), [pt],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.0)))
    # 教訓
    big_banner(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.65),
               [[("「自分でやってから途中で依頼」は、", 20, WHITE, False, SANS),
                 ("効率的とは限らない", 22, WHITE, True, SANS)]], fill=GREEN)

# ---- CASE 03 ----
def s10_case03_door():
    s = slide()
    case_band(s, "CASE 03", ["外車の修理費対応の中で", "評価損の請求漏れが見つかった事例"],
              [[("修理費が出れば終わり —— とは限らない。", LEAD, TEXT_SUB, False, SANS),
                ("評価損", LEAD, GREEN_DARK, True, SANS),
                ("という論点が生じる", LEAD, TEXT_SUB, False, SANS)]])
    kw_card(s, Inches(0.7), Inches(4.4), Inches(12.0), Inches(2.4),
            [("対象車両は", 22, TEXT, False, SANS),
             ("外車・購入1年未満", 22, GREEN, True, SANS)],
            body="相手方との調整がうまく進まず、保険代理店様を通じて弁護士に相談が来た事案。",
            body_size=22)

def s11_case03_find():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Case 03")
    heading(s, Inches(0.7), Inches(1.1), Inches(8.3), "修理費だけでは足りなかった")
    add_picture_cover(s, "評価損", Inches(9.3), Inches(0.6), Inches(3.4), Inches(3.4))
    rows = [
        em("当初", normal_after="：修理見積を出すだけの予定だった", size=22),
        [("弁護士の確認で ", 22, TEXT, False, SANS), ("評価損", 24, GREEN, True, SANS), (" も請求できると判明", 22, TEXT, False, SANS)],
    ]
    y = Inches(2.4)
    for r in rows:
        rect(s, Inches(0.7), y, Inches(8.0), Inches(0.95), fill=GREEN_LT)
        txt(s, Inches(1.0), y, Inches(7.5), Inches(0.95), [r],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.1)))
    big_banner(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(1.25),
               [[("外車 × 購入1年未満", 30, WHITE, True, SERIF),
                 (" → 事故歴で将来の市場価値が下がる", 24, WHITE, False, SANS)]])

def s12_case03_check():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Case 03")
    heading(s, Inches(0.7), Inches(1.1), Inches(12), "評価損を検討すべきケース")
    cards = [
        ([("外車・輸入車", H3, GREEN_DARK, True, SANS)], "事故歴による市場価値の低下が国産車より大きい"),
        ([("購入から年数が浅い", 24, GREEN_DARK, True, SANS)], "購入後1〜2年以内は認められる可能性が特に高い"),
        ([("修理費が高額", H3, GREEN_DARK, True, SANS)], "車両価格の一定割合を超える損傷は根拠になりやすい"),
    ]
    cw = Inches(3.85); ch = Inches(2.6); gx = Inches(0.22)
    x0 = Inches(0.7); y0 = Inches(2.4)
    for i, (head, body) in enumerate(cards):
        x = Emu(int(x0) + i * (int(cw) + int(gx)))
        kw_card(s, x, y0, cw, ch, head, body=body, body_size=18)
    big_banner(s, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.1),
               [[("保険代理店様へ：", 22, GREEN_MID, True, SANS)],
                [("外車・高額車両・購入間もない車両は、", 24, WHITE, False, SANS),
                 ("見積もり段階から速やかにご相談を", 26, WHITE, True, SERIF)]], fill=GREEN_DARK)

# ---- 結論 ----
def s13_summary_table():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Conclusion")
    heading(s, Inches(0.7), Inches(1.1), Inches(9), "3つに共通すること", size=H3 + 6)
    add_picture_cover(s, "家族", Inches(11.0), Inches(0.5), Inches(1.9), Inches(2.0))
    rows = [
        ["事例", "途中介入で生じた課題", "初動介入であれば"],
        ["① 治療費打切り", "請求が後ろ倒し・社内調整やり直し・長期化", "当初から全請求項目を整理・円滑な折衝"],
        ["② 後遺障害", "一度不認定 → やり直しで時間と労力", "認定を見据えた診断書・資料整備"],
        ["③ 外車・評価損", "後から評価損を追加して複雑化", "当初から請求項目に組み込み整合的に交渉"],
    ]
    tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(2.6),
                             Inches(12.0), Inches(3.2)).table
    tbl.columns[0].width = Inches(2.9)
    tbl.columns[1].width = Inches(4.55)
    tbl.columns[2].width = Inches(4.55)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.line_spacing = 1.15
            r = p.add_run(); r.text = val
            if ri == 0:
                r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
                _set_ea_font(r, SERIF)
                cell.fill.solid(); cell.fill.fore_color.rgb = GREEN
            elif ci == 0:
                r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = GREEN_DARK
                _set_ea_font(r, SERIF)
                cell.fill.solid(); cell.fill.fore_color.rgb = GREEN_LT
            else:
                r.font.size = Pt(18); r.font.color.rgb = TEXT
                _set_ea_font(r, SANS)
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    big_banner(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(1.0),
               [[("途中でも解決はできる。が、", 24, WHITE, False, SANS),
                 ("初動の方が早くスムーズ", 28, WHITE, True, SERIF)]])

def s14_benefits():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Conclusion")
    heading(s, Inches(0.7), Inches(1.1), Inches(12), "最初から入ると、何が違うか")
    claims = [
        ([("請求項目の漏れ", H3, GREEN_DARK, True, SANS), ("を防げる", H3, GREEN_DARK, True, SANS)],
         "交通事故の全請求項目を把握し、初期から漏れのない請求リストを作成できる"),
        ([("後遺障害認定", H3, GREEN_DARK, True, SANS), ("を見据えられる", H3, GREEN_DARK, True, SANS)],
         "認定基準を理解した弁護士が主治医と連携し、診断書・検査・資料を最適化"),
        ([("評価損", H3, GREEN_DARK, True, SANS), ("など見落としを早期に拾える", H3, GREEN_DARK, True, SANS)],
         "代車費用・格落ち損害など見落とされやすい論点も初期から一体的に請求"),
    ]
    y = Inches(2.3)
    for head, body in claims:
        rect(s, Inches(0.7), y, Inches(12.0), Inches(1.4), fill=GREEN_LT)
        # ✓バッジ
        rect(s, Inches(1.0), Emu(int(y) + int(Inches(0.45))), Inches(0.55), Inches(0.55),
             fill=GREEN, shape=MSO_SHAPE.OVAL)
        txt(s, Inches(1.0), Emu(int(y) + int(Inches(0.45))), Inches(0.55), Inches(0.55),
            [P("✓", 24, WHITE, True, SANS)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
        txt(s, Inches(1.85), Emu(int(y) + int(Inches(0.2))), Inches(10.6), Inches(1.0),
            [head, P(body, 18, TEXT, False, SANS)],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2, space_after=Pt(2))
        y = Emu(int(y) + int(Inches(1.6)))

def s15_tokuyaku():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "Message")
    heading(s, Inches(0.7), Inches(1.1), Inches(8), "弁護士費用特約の論理")
    add_picture_cover(s, "相談", Inches(9.3), Inches(0.6), Inches(3.4), Inches(3.2))
    # 強調2行
    rect(s, Inches(0.7), Inches(2.5), Inches(8.0), Inches(1.1), fill=GREEN_LT)
    txt(s, Inches(1.0), Inches(2.5), Inches(7.5), Inches(1.1),
        [[("特約がある限り、", 22, TEXT, False, SANS),
          ("依頼者の自己負担なし", 24, GREEN, True, SANS)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=Pt(0))
    rect(s, Inches(0.7), Inches(3.75), Inches(8.0), Inches(1.1), fill=GREEN_LT)
    txt(s, Inches(1.0), Inches(3.75), Inches(7.5), Inches(1.1),
        [[("だから ", 22, TEXT, False, SANS),
          ("早く使うほどメリットが大きい", 24, GREEN, True, SANS)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=Pt(0))
    big_banner(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(1.1),
               [[("紹介元（保険代理店様）への ", 22, WHITE, False, SANS),
                 ("顧客からの信頼向上", 26, WHITE, True, SERIF),
                 (" にもつながる", 22, WHITE, False, SANS)]], fill=GREEN_DARK)

def s16_keymessage():
    s = slide()
    rect(s, 0, 0, SW, SH, fill=GREEN_DARK)
    txt(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.6),
        [[("「揉めてから」", 44, GREEN_MID, True, SERIF), ("ではなく", 36, WHITE, False, SANS)],
         [("「揉める前」", 54, WHITE, True, SERIF), ("に弁護士へ", 44, WHITE, True, SERIF)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25, space_after=Pt(0))
    txt(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.8),
        [P("依頼者・紹介元・相手方保険会社、いずれにとっても合理的な選択", LEAD, GREEN_MID, False, SANS)],
        align=PP_ALIGN.CENTER, line_spacing=1.2, space_after=Pt(0))

def s17_cta():
    s = slide()
    add_picture_cover(s, "cta", 0, 0, SW, SH)
    # 左に濃緑ボックス（視認性）
    rect(s, 0, 0, Inches(7.6), SH, fill=GREEN_DARK)
    section_label(s, Inches(0.7), Inches(0.9), "Contact", color=GREEN_MID)
    txt(s, Inches(0.7), Inches(1.6), Inches(6.6), Inches(1.8),
        [P("保険代理店様との", 34, WHITE, True, SERIF),
         P("連携をお待ちしています", 34, WHITE, True, SERIF)],
        line_spacing=1.2, space_after=Pt(0))
    txt(s, Inches(0.7), Inches(3.6), Inches(6.6), Inches(1.2),
        [P("初回相談は無料。対面・オンライン（Zoom等）どちらでも対応します。", 20, GREEN_LT, False, SANS)],
        line_spacing=1.35, space_after=Pt(0))
    badges = ["初回相談無料", "オンライン全国対応", "弁護士費用特約案件 多数"]
    y = Inches(5.1)
    for b in badges:
        rect(s, Inches(0.7), y, Inches(4.6), Inches(0.6), fill=GREEN_LT)
        txt(s, Inches(0.7), y, Inches(4.6), Inches(0.6),
            [P(b, 20, GREEN_DARK, True, SANS)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(0.72)))

def s18_about():
    s = slide()
    section_label(s, Inches(0.7), Inches(0.6), "About")
    heading(s, Inches(0.7), Inches(1.1), Inches(12), "事務所概要")
    rows_left = [
        ("所在地", "〒104-0061 東京都中央区銀座8丁目17-5"),
        ("事務所名", "弁護士法人 相生綜合法律事務所"),
        ("代表弁護士", "福島 駿太（第一東京弁護士会 登録）"),
    ]
    rows_right = [
        ("取扱分野", "交通事故、離婚・相続、労働、債務整理、企業法務 ほか"),
        ("対応エリア", "オンラインにて全国対応"),
        ("初回相談", "無料（対面・オンライン対応）"),
    ]
    txt(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(3.0),
        [[(lbl + "：", 20, GREEN, True, SANS), (val, 20, TEXT, False, SANS)] for lbl, val in rows_left],
        line_spacing=1.4, space_after=Pt(16))
    txt(s, Inches(7.0), Inches(2.4), Inches(5.9), Inches(3.0),
        [[(lbl + "：", 20, GREEN, True, SANS), (val, 20, TEXT, False, SANS)] for lbl, val in rows_right],
        line_spacing=1.4, space_after=Pt(16))
    # 免責（小・帯）
    rect(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.3), fill=GREEN_LT)
    txt(s, Inches(1.0), Inches(5.85), Inches(11.4), Inches(1.1),
        [P("本事例は実際の案件をもとに整理・匿名化したもので、個別の結果を保証するものではありません。", 16, TEXT_SUB, False, SANS),
         P("記載内容は2026年4月時点の情報です。具体的なご相談は弁護士までお願いいたします。", 16, TEXT_SUB, False, SANS)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3, space_after=Pt(4))

# ================================================================ ビルド・検査・保存

def build_all():
    s01_cover(); s02_intro(); s03_conclusion_first()
    s04_case01_door(); s05_case01_leak(); s06_case01_lesson()
    s07_case02_door(); s08_case02_flow(); s09_case02_lesson()
    s10_case03_door(); s11_case03_find(); s12_case03_check()
    s13_summary_table(); s14_benefits(); s15_tokuyaku()
    s16_keymessage(); s17_cta(); s18_about()

def verify():
    """はみ出し・写真枚数・文字量をチェック。"""
    issues = []
    pic_total = 0
    SWp, SHp = float(SW), float(SH)
    for si, sl in enumerate(prs.slides, 1):
        chars = 0
        for sh in sl.shapes:
            if sh.shape_type == 13:
                pic_total += 1
            # 枠はみ出し（左上基準・図形のみ）
            try:
                if sh.left is not None and sh.top is not None and sh.width and sh.height:
                    if sh.left < -1000 or sh.top < -1000:
                        pass  # cover画像は意図的に負座標になる場合があるので無視
                    if int(sh.left) + int(sh.width) > int(SW) + Emu(0.05 * 914400):
                        issues.append((si, "右はみ出し", sh.shape_type))
                    if int(sh.top) + int(sh.height) > int(SH) + Emu(0.05 * 914400):
                        issues.append((si, "下はみ出し", sh.shape_type))
            except Exception:
                pass
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        chars += len(r.text)
            if sh.has_table:
                for row in sh.table.rows:
                    for c in row.cells:
                        for p in c.text_frame.paragraphs:
                            for r in p.runs:
                                chars += len(r.text)
        if chars > 320:
            issues.append((si, f"文字過多 {chars}字", None))
    print(f"[INFO] 写真={pic_total}枚 / スライド={len(prs.slides._sldIdLst)}枚")
    if issues:
        for it in issues:
            print("  注意:", it)
    else:
        print("[OK] はみ出し・文字量チェック問題なし")

if __name__ == "__main__":
    build_all()
    verify()
    try:
        prs.save(OUT)
    except Exception as e:
        raise SystemExit(f"[NG] 保存失敗（ファイルが開かれている可能性）: {e}")
    print(f"[OK] 保存しました: {OUT}  スライド数={len(prs.slides._sldIdLst)}")
