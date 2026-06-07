# -*- coding: utf-8 -*-
"""
特別レポートHTML → 編集可能PPTX（16:9・最小フォント16pt）

- 元: 20260416_特別レポート（コンテンツ）.html（A4縦・全12ページ）
- 出力: 20260416_特別レポート（16-9）.pptx
- 方針: python-pptxでネイティブ図形＋テキストボックスとして1から構築（完全編集可能）
- 本文は元HTMLの原文を1字1句そのまま使用（要約・言い換えをしない）
- 最小フォント16pt（フッター・ラベル含め16未満を作らない。生成後に自動検査）
- 画像は元HTMLから抽出した実写真（_html_img/）を alt 文言で対応付けて配置
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 定数

BASE = r"g:\共有ドライブ\マーケティング部\23_保険代理店開拓"
IMG  = os.path.join(BASE, "_html_img")
OUT  = os.path.join(BASE, "20260416_特別レポート（16-9）.pptx")

# 16:9 スライド寸法
SW = Inches(13.333)
SH = Inches(7.5)

# HTMLの配色を再現
GREEN      = RGBColor(0x2B, 0x64, 0x44)   # --green
GREEN_DARK = RGBColor(0x1E, 0x4D, 0x33)   # --green-dark
GREEN_LT   = RGBColor(0xE8, 0xF2, 0xEB)   # --green-light
GREEN_MID  = RGBColor(0xC8, 0xDF, 0xD0)   # --green-mid
SAGE       = RGBColor(0xD4, 0xE4, 0xD8)
TEXT       = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SUB   = RGBColor(0x55, 0x55, 0x55)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GOLD       = RGBColor(0xB8, 0x96, 0x4A)

SANS  = "Noto Sans JP"     # 本文（システムにVF版インストール済み）
SERIF = "Noto Serif JP"    # 見出し

MIN_PT = 16  # 最小フォント

# 画像（alt→ファイル）
IMGS = {
    "cover":   os.path.join(IMG, "cover_1_cover-right.jpg"),   # 緑の木漏れ日
    "依頼者":  os.path.join(IMG, "01_依頼者.jpg"),
    "書類対応": os.path.join(IMG, "02_書類対応.jpg"),
    "法律書籍": os.path.join(IMG, "03_法律書籍.jpg"),
    "評価損":  os.path.join(IMG, "04_評価損.jpg"),
    "家族":    os.path.join(IMG, "05_家族.jpg"),
    "相談":    os.path.join(IMG, "06_相談.jpg"),
    "相談室":  os.path.join(IMG, "07_相談室.jpg"),
    "弁護士":  os.path.join(IMG, "08_弁護士.jpg"),
    "cta":     os.path.join(IMG, "bg_cta-bg-img.jpg"),
}

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]  # 空白レイアウト

# ---------------------------------------------------------------- ヘルパ

def _set_ea_font(run, name):
    """日本語（East Asian）フォントを確実に効かせる。latin/ea 両方に同名を設定。"""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE,
         shadow=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    if not shadow:
        # 影を明示的にオフ（python-pptxはテーマ影を引き継ぐことがある）
        pass
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.25, space_after=Pt(6), wrap=True):
    """runs: list of paragraphs; 各paragraph = list of (text, size, color, bold, font)
       size は最小16ptを保証（呼び出し側でも守るが二重で担保）。"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_after = space_after
        p.space_before = Pt(0)
        for (t, size, color, bold, font) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(max(size, MIN_PT))
            r.font.color.rgb = color
            r.font.bold = bold
            _set_ea_font(r, font)
    return tb

def P(text, size=16, color=TEXT, bold=False, font=SANS):
    """1ラン1段落の糖衣。"""
    return [(text, size, color, bold, font)]

def pic_cover(s, path, x, y, w, h):
    """指定枠を完全に覆うように画像を中央クロップ配置（背景写真用）。"""
    from PIL import Image  # 使えないので自前計算に切替（下のフォールバック）
    return None

def add_picture_cover(s, path, x, y, w, h):
    """枠(w,h)をアスペクト比維持で覆い、はみ出しはクロップ。PILなしでpx寸法取得。"""
    iw, ih = _img_px(path)
    if not iw:
        s.shapes.add_picture(path, x, y, width=w, height=h)
        return
    fw, fh = float(w), float(h)
    fr = fw / fh
    ir = iw / ih
    if ir > fr:
        # 画像が横長 → 高さを合わせて横をクロップ
        new_h = fh
        new_w = fh * ir
    else:
        new_w = fw
        new_h = fw / ir
    px = x - (new_w - fw) / 2
    py = y - (new_h - fh) / 2
    pic = s.shapes.add_picture(path, int(px), int(py), int(new_w), int(new_h))
    pic.crop_left = max(0, (new_w - fw) / 2 / new_w)
    pic.crop_right = max(0, (new_w - fw) / 2 / new_w)
    pic.crop_top = max(0, (new_h - fh) / 2 / new_h)
    pic.crop_bottom = max(0, (new_h - fh) / 2 / new_h)
    # 配置は枠内に戻す
    pic.left = int(x); pic.top = int(y); pic.width = int(w); pic.height = int(h)
    return pic

def _img_px(path):
    """JPEG/PNGのpx寸法を外部ライブラリなしで読む。"""
    import struct
    try:
        with open(path, "rb") as f:
            data = f.read()
        if data[:8] == b"\x89PNG\r\n\x1a\n":
            return struct.unpack(">II", data[16:24])
        if data[:2] == b"\xff\xd8":
            i = 2; n = len(data)
            while i < n:
                if data[i] != 0xFF:
                    i += 1; continue
                mk = data[i+1]
                if mk in (0xC0,0xC1,0xC2,0xC3,0xC5,0xC6,0xC7,0xC9,0xCA,0xCB,0xCD,0xCE,0xCF):
                    h = struct.unpack(">H", data[i+5:i+7])[0]
                    w = struct.unpack(">H", data[i+7:i+9])[0]
                    return w, h
                i += 2 + struct.unpack(">H", data[i+2:i+4])[0]
    except Exception:
        pass
    return None, None

def footer(s, page_no):
    """全ページ共通フッター。16pt。"""
    txt(s, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
        [P("弁護士法人 相生綜合法律事務所", 16, TEXT_SUB, False, SANS)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
    txt(s, Inches(12.0), Inches(7.05), Inches(0.83), Inches(0.35),
        [P(f"P. {page_no:02d}", 16, TEXT_SUB, False, SANS)],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))

def section_label(s, x, y, text, color=GREEN):
    """英字の小見出しラベル（HTMLの.section-label相当）。16pt。"""
    txt(s, x, y, Inches(6), Inches(0.4),
        [P(text.upper(), 16, color, True, SANS)],
        line_spacing=1.0, space_after=Pt(0))

def heading(s, x, y, w, text, size=30, color=GREEN_DARK):
    """セクション見出し（Serif）。"""
    return txt(s, x, y, w, Inches(1.0),
               [P(text, size, color, True, SERIF)],
               line_spacing=1.1, space_after=Pt(0))

def h3(s, x, y, w, text, size=20, color=GREEN_DARK):
    return txt(s, x, y, w, Inches(0.6),
               [P(text, size, color, True, SERIF)],
               line_spacing=1.15, space_after=Pt(0))

def num_badge(s, x, y, n, dia=Inches(0.5), fill=GREEN, fg=WHITE):
    sp = rect(s, x, y, dia, dia, fill=fill, shape=MSO_SHAPE.OVAL)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = fg
    _set_ea_font(r, SANS)
    sp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return sp

def case_band(s, case_no, title_lines):
    """CASE扉の濃緑グラデ帯（上部）。CASE番号＋タイトル。"""
    band = rect(s, 0, 0, SW, Inches(1.85), fill=GREEN_DARK)
    # 右側にうっすら明るい緑のアクセント帯
    rect(s, Inches(10.8), 0, Inches(2.6), Inches(1.85), fill=GREEN)
    txt(s, Inches(0.6), Inches(0.32), Inches(4), Inches(0.45),
        [P(case_no, 18, GREEN_MID, True, SANS)], line_spacing=1.0, space_after=Pt(0))
    txt(s, Inches(0.6), Inches(0.7), Inches(12.1), Inches(1.0),
        [P(l, 26, WHITE, True, SERIF) for l in title_lines],
        line_spacing=1.12, space_after=Pt(0))

# ================================================================ スライド構築

def s01_cover():
    s = slide()
    # 左：白地テキスト / 右：木漏れ日写真
    rect(s, 0, 0, Inches(8.0), SH, fill=WHITE)
    add_picture_cover(s, IMGS["cover"], Inches(8.0), 0, Inches(5.333), SH)
    # 緑のオーバーレイ（写真上に薄く）
    ov = rect(s, Inches(8.0), 0, Inches(5.333), SH, fill=GREEN_DARK)
    ov.fill.fore_color.rgb = GREEN_DARK
    ov.fill.transparency = 0  # python-pptxでは透過は別途。視認性のため左境界に細い緑バー
    # 透過が効かないので、写真はそのまま見せ、左テキスト側に集約する。上のovを消す:
    ov._element.getparent().remove(ov._element)
    # 左テキスト
    txt(s, Inches(0.7), Inches(0.7), Inches(6.5), Inches(0.5),
        [P("SPECIAL REPORT", 17, GREEN, True, SANS)], line_spacing=1.0, space_after=Pt(0))
    # タグ
    tag = rect(s, Inches(0.7), Inches(1.25), Inches(3.2), Inches(0.5), fill=GREEN_LT)
    txt(s, Inches(0.7), Inches(1.25), Inches(3.2), Inches(0.5),
        [P("交通事故 弁護士費用特約", 16, GREEN_DARK, True, SANS)],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
    # タイトル
    txt(s, Inches(0.7), Inches(2.0), Inches(7.0), Inches(1.8),
        [P("交通事故案件における", 38, GREEN_DARK, True, SERIF),
         P("弁護士早期介入の重要性", 38, GREEN_DARK, True, SERIF)],
        line_spacing=1.15, space_after=Pt(0))
    txt(s, Inches(0.7), Inches(3.7), Inches(7.0), Inches(0.6),
        [P("3つの事例から考える「初動から弁護士を入れる」ことの意義", 18, TEXT_SUB, False, SANS)],
        line_spacing=1.2, space_after=Pt(0))
    # リード
    txt(s, Inches(0.7), Inches(4.4), Inches(7.0), Inches(0.9),
        [P("本レポートは、弁護士法人相生綜合法律事務所が実際に取り扱った交通事故案件をもとに、保険代理店様へ向けて作成した特別レポートです。",
           16, TEXT, False, SANS)],
        line_spacing=1.4, space_after=Pt(0))
    # 統計ボックス3つ
    stats = [("3", "実際の事例をもとに解説"), ("0円", "弁護士費用特約\n依頼者負担なし"), ("初回", "無料相談\n対面・オンライン")]
    bx = Inches(0.7); bw = Inches(2.2); gap = Inches(0.15)
    for i, (big, small) in enumerate(stats):
        x = Emu(int(bx) + i * (int(bw) + int(gap)))
        rect(s, x, Inches(5.45), bw, Inches(1.1), fill=GREEN)
        txt(s, x, Inches(5.5), bw, Inches(0.55),
            [P(big, 28, WHITE, True, SERIF)], align=PP_ALIGN.CENTER,
            line_spacing=1.0, space_after=Pt(0))
        txt(s, x, Inches(6.05), bw, Inches(0.5),
            [P(small.replace("\n", " "), 16, WHITE, False, SANS)], align=PP_ALIGN.CENTER,
            line_spacing=1.05, space_after=Pt(0))
    # フッター（表紙用）
    txt(s, Inches(0.7), Inches(6.85), Inches(7.0), Inches(0.6),
        [[("弁護士法人 相生綜合法律事務所　", 16, GREEN_DARK, True, SERIF),
          ("代表弁護士 福島駿太（第一東京弁護士会）　2026年4月", 16, TEXT_SUB, False, SANS)]],
        line_spacing=1.1, space_after=Pt(0))

def s02_intro1():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Introduction")
    heading(s, Inches(0.6), Inches(0.95), Inches(8), "はじめに", size=32)
    # 右に写真（依頼者＝手をつなぐ）
    add_picture_cover(s, IMGS["依頼者"], Inches(9.4), Inches(0.6), Inches(3.3), Inches(2.5))
    txt(s, Inches(0.6), Inches(2.1), Inches(8.4), Inches(4.5),
        [P("保険代理店の方々から、「弁護士にはいつ相談すればよいのか」「揉めていなくても依頼してよいものか」といったご質問をいただくことがあります。", 18, TEXT, False, SANS),
         P("結論から申し上げると、弁護士費用特約が使える案件であれば、事故後できるだけ早い段階で弁護士につないでいただくことをお勧めします。", 18, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(14))
    # 強調ボックス
    rect(s, Inches(9.4), Inches(3.4), Inches(3.3), Inches(3.0), fill=GREEN_LT)
    txt(s, Inches(9.65), Inches(3.65), Inches(2.8), Inches(2.6),
        [P("POINT", 16, GREEN, True, SANS),
         P("弁護士費用特約が使えるなら、事故後できるだけ早い段階で弁護士へ。", 18, GREEN_DARK, True, SERIF)],
        line_spacing=1.4, space_after=Pt(10))
    footer(s, 2)

def s03_intro2():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Introduction / 続き")
    heading(s, Inches(0.6), Inches(0.95), Inches(11), "はじめに（続き）", size=28)
    txt(s, Inches(0.6), Inches(1.95), Inches(12.1), Inches(4.6),
        [P("本レポートでは、実際に当事務所が対応した3つの事例を通じて、「途中から弁護士が入る」ことと「最初から弁護士が入る」ことの違いを具体的に示します。いずれの事例も、最終的には適切な解決に至っていますが、初動から弁護士が関与していれば、より早く、よりスムーズに進められたことが共通しています。", 18, TEXT, False, SANS),
         P("事故後の初動に関わる方が、「まだ揉めていないから様子を見よう」ではなく、「弁護士費用特約があるなら、まず弁護士につないでおこう」という発想で動いていただくことに、依頼者にとって大きな意味があります。本レポートがその判断の一助となれば幸いです。", 18, TEXT, False, SANS)],
        line_spacing=1.55, space_after=Pt(16))
    footer(s, 3)

def s04_toc():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Contents")
    heading(s, Inches(0.6), Inches(0.95), Inches(11), "本レポートの内容", size=32)
    items = [
        ("01", "治療費打切り後に請求漏れが見つかった事例", GREEN),
        ("02", "後遺障害診断書の記載を修正したことで、後遺障害認定が得られた事例", GREEN),
        ("03", "外車の修理費対応の中で評価損の請求漏れが見つかった事例", GREEN),
        ("結", "3つの事例に共通していること —— 初動から弁護士が関与することの利点", GREEN_DARK),
    ]
    y = Inches(2.2)
    for num, label, col in items:
        rect(s, Inches(0.6), y, Inches(0.9), Inches(0.9), fill=col, shape=MSO_SHAPE.OVAL)
        txt(s, Inches(0.6), y, Inches(0.9), Inches(0.9),
            [P(num, 26, WHITE, True, SERIF)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
        txt(s, Inches(1.8), y, Inches(10.8), Inches(0.9),
            [P(label, 19, TEXT, False, SANS)], anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.25, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.12)))
    footer(s, 4)

# ---- CASE 01 ----
def s05_case01_a():
    s = slide()
    case_band(s, "CASE 01", ["治療費打切り後に請求漏れが見つかった事例"])
    txt(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.3),
        [P("むち打ち等で通院中、相手方保険会社から「治療費の打切り」を告げられてから初めて相談に来られるケースは少なくありません。しかし、その段階で内容を確認すると、本来請求できるはずの項目が相当程度漏れていることがあります。", 17, TEXT_SUB, False, SANS)],
        line_spacing=1.45, space_after=Pt(0))
    num_badge(s, Inches(0.6), Inches(3.6), 1)
    h3(s, Inches(1.3), Inches(3.62), Inches(11), "相談が多い典型的な場面")
    txt(s, Inches(1.3), Inches(4.2), Inches(11.4), Inches(2.6),
        [P("交通事故の人身案件では、保険代理店様などからご紹介をいただくことが多いです。その中でも特に多いのが、お客様がむち打ち等で整骨院や整形外科に通っていて、相手方保険会社の一括対応で治療費が支払われていたものの、「今月末で治療費の支払が打切りになります」と言われ、「この後どうしたらよいでしょうか」と相談に来られるケースです。", 17, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(0))
    footer(s, 5)

def s06_case01_b():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 01 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 2)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "実際に見つかる請求漏れの内容")
    txt(s, Inches(0.6), Inches(1.7), Inches(7.4), Inches(4.6),
        [P("このような場合、実際に内容を確認してみると、相手方保険会社との間でやり取りされているのが、治療費の直接支払に関する話だけであることが少なくありません。入通院慰謝料の話が全くされていないこともありますし、交通費の請求がされていないこともあります。また、自営業者や会社経営者の方であれば、本来は休業損害や営業損害を検討すべきであるにもかかわらず、その点の請求が全くされていないこともあります。つまり、本来請求できるはずの項目が、相当程度漏れてしまっていることがあるのです。", 17, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(0))
    # 右：ハイライトボックス（4項目）
    rect(s, Inches(8.3), Inches(1.7), Inches(4.4), Inches(4.9), fill=GREEN_LT)
    txt(s, Inches(8.6), Inches(1.95), Inches(3.9), Inches(0.6),
        [P("請求漏れとして見つかりやすい主な項目", 18, GREEN_DARK, True, SERIF)],
        line_spacing=1.2, space_after=Pt(0))
    bullets = [
        "入通院慰謝料（通院期間に応じた精神的損害の補償）",
        "交通費（通院時の交通機関利用料、駐車場代等）",
        "休業損害（事故による収入減少の補償）",
        "営業損害（自営業者・経営者に特有の損害項目）",
    ]
    txt(s, Inches(8.6), Inches(2.75), Inches(3.9), Inches(3.7),
        [[("●  ", 16, GREEN, True, SANS), (b, 16, TEXT, False, SANS)] for b in bullets],
        line_spacing=1.35, space_after=Pt(12))
    footer(s, 6)

def s07_case01_c():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 01 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 3)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "弁護士介入後の対応")
    txt(s, Inches(0.6), Inches(1.7), Inches(8.3), Inches(4.8),
        [P("そこで弁護士が介入すると、まずは整形外科や整骨院側に、なお治療継続の必要があることを示す診断書や意見書等を作成してもらい、それを相手方保険会社に提出して、必要な範囲で治療期間の延長を求めていくことになります。そして、それと並行して、それまで見落とされていた交通費、休業損害、営業損害、入通院慰謝料などの請求項目を整理し、まとめて請求していくことになります。こうした対応をすることで、結果として適切な解決に至ることが多いです。", 18, TEXT, False, SANS)],
        line_spacing=1.55, space_after=Pt(0))
    add_picture_cover(s, IMGS["書類対応"], Inches(9.2), Inches(1.7), Inches(3.5), Inches(2.5))
    footer(s, 7)

def s08_case01_d():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 01 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 4)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "この事例から見える問題点")
    txt(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.8),
        [P("もっとも、このような案件に接していていつも感じるのは、やはりもっと早い段階で相談していただいた方がよいということです。気づくのが遅れれば、それだけ依頼者のもとに入るべきお金が入ってくるのも遅れますし、請求が後ろ倒しになれば支払時期も後ろ倒しになります。その意味で、依頼者にとっては明らかに不利益です。", 18, TEXT, False, SANS),
         P("また、保険会社側から見ても、いったん「この内容で終了する」という前提で社内決裁や方針整理が進んだ後に、途中から弁護士が入って治療期間の延長や追加請求を行うことになると、再度社内調整が必要になります。その結果、案件全体が長期化しやすくなります。依頼者のためにも、相手方保険会社とのやり取りを円滑に進めるためにも、最初から請求項目を整理して進めた方がよいのです。", 18, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(14))
    footer(s, 8)

def s09_case01_e():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 01 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 5)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "依頼者の受け止め方と初動対応の重要性")
    txt(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(3.2),
        [P("さらに、依頼者の感覚としては、「相手方保険会社も十分に教えてくれなかった」「保険代理店もそこまで案内してくれなかった」という受け止め方になりやすいです。しかし、最初から弁護士が入っていれば、相手方保険会社としても「弁護士が関与している案件である」という前提で対応しますので、当初から請求項目を意識した整理がしやすくなります。結果として、後から修正を重ねるよりも、最初から弁護士に入ってもらった方が、むしろシンプルで早いということが多いです。", 18, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(0))
    # プルクオート
    rect(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.5), fill=GREEN)
    txt(s, Inches(1.0), Inches(5.05), Inches(11.3), Inches(1.4),
        [P("後から修正を重ねるよりも、最初から弁護士に入ってもらった方が、むしろシンプルで早いということが多い。", 20, WHITE, True, SERIF)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.35, space_after=Pt(0))
    footer(s, 9)

# ---- CASE 02 ----
def s10_case02_a():
    s = slide()
    case_band(s, "CASE 02", ["後遺障害診断書の記載を修正したことで", "後遺障害認定が得られた事例"])
    txt(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.3),
        [P("後遺障害認定は、診断書を提出すれば自動的に得られるものではありません。認定基準との関係で、どの症状をどのように記載するかが結果を左右します。この事例は、弁護士の関与によって不認定が覆った実例です。", 17, TEXT_SUB, False, SANS)],
        line_spacing=1.45, space_after=Pt(0))
    num_badge(s, Inches(0.6), Inches(3.55), 1)
    h3(s, Inches(1.3), Inches(3.57), Inches(11), "相談時点の状況")
    txt(s, Inches(1.3), Inches(4.15), Inches(11.4), Inches(2.6),
        [P("次に、後遺障害診断書の記載を修正したことで、最終的に後遺障害認定が得られた事例があります。この件は、保険代理店からお客様をご紹介いただいたものでした。相談時点では、すでに後遺障害診断書を取得し、後遺障害認定機関に提出済みという状況でした。", 17, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(0))
    footer(s, 10)

def s11_case02_b():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 02 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 2)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "一度は不認定となった経緯")
    txt(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.6),
        [P("もっとも、その段階ではすでに提出が済んでいたため、まずは結果を待つしかありませんでした。そして、実際に結果が出てみると、不認定という判断になっていました。そこで改めて弁護士が介入し、診断書を作成した整形外科医と面談又は協議を行い、どのような症状をどのように記載すべきか、どの点が認定上重要になるのかを整理したうえで、診断書の内容を適切に修正してもらいました。その修正版を再度提出したところ、最終的には後遺障害認定を得ることができました。", 17, TEXT, False, SANS)],
        line_spacing=1.45, space_after=Pt(0))
    num_badge(s, Inches(0.6), Inches(4.5), 3)
    h3(s, Inches(1.3), Inches(4.52), Inches(11), "この事例から分かること")
    txt(s, Inches(1.3), Inches(5.1), Inches(11.4), Inches(1.9),
        [P("この事例から分かるのは、自分でできるところまでやってから途中で弁護士に依頼するという進め方は、一見すると効率的に見えても、実際にはそうとは限らないということです。いったん不認定という結果が出てしまうと、そこから再度やり直す必要があり、時間も労力も余分にかかります。最初から弁護士が関与していれば、後遺障害認定を見据えた形で診断書の記載内容や必要資料を整理しながら進めることができたはずです。", 16, TEXT, False, SANS)],
        line_spacing=1.4, space_after=Pt(0))
    footer(s, 11)

def s12_case02_c():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 02 / 続き")
    # ハイライトボックス
    rect(s, Inches(0.6), Inches(1.0), Inches(12.1), Inches(1.5), fill=GREEN_LT)
    txt(s, Inches(0.95), Inches(1.1), Inches(11.4), Inches(1.3),
        [P("後遺障害認定は、診断書の記載内容が結果を直接左右します。最初から弁護士が関与することで、認定基準を踏まえた準備が可能となり、不認定からのやり直しを防ぐことができます。", 18, GREEN_DARK, True, SERIF)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4, space_after=Pt(0))
    h3(s, Inches(0.6), Inches(2.9), Inches(8.3), "後遺障害の場面で重要な視点")
    txt(s, Inches(0.6), Inches(3.55), Inches(8.3), Inches(3.2),
        [P("特に後遺障害の場面では、単に診断書を出せばよいというものではなく、認定基準との関係で、どの症状をどのように記載するかが非常に重要です。そこは一般の方が当然に理解できる領域ではありませんし、医師側も認定実務を前提とした記載のポイントまで十分に把握しているとは限りません。だからこそ、認定を見据えた段階で弁護士が関与していることに意味があります。", 17, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(0))
    add_picture_cover(s, IMGS["法律書籍"], Inches(9.2), Inches(2.9), Inches(3.5), Inches(3.4))
    footer(s, 12)

def s13_case02_d():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 02 / 続き")
    h3(s, Inches(0.6), Inches(1.0), Inches(11), "初動から弁護士が関与する意義", size=22)
    # 緑見出しボックス
    rect(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.7), fill=GREEN)
    txt(s, Inches(0.95), Inches(1.8), Inches(11.4), Inches(0.7),
        [P("後遺障害認定において弁護士が果たす役割", 20, WHITE, True, SERIF)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
    points = [
        "認定基準を踏まえた診断書記載のポイントを医師に伝える",
        "認定に必要な検査・画像資料を初期段階から整備する",
        "不認定の場合の異議申立を視野に入れた準備を並行して進める",
    ]
    y = Inches(2.85)
    for i, pt in enumerate(points, 1):
        rect(s, Inches(0.6), y, Inches(12.1), Inches(1.05), fill=GREEN_LT)
        num_badge(s, Inches(0.9), Emu(int(y)+int(Inches(0.27))), i, dia=Inches(0.5))
        txt(s, Inches(1.7), y, Inches(10.8), Inches(1.05),
            [P(pt, 19, GREEN_DARK, True, SANS)], anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.25, space_after=Pt(0))
        y = Emu(int(y) + int(Inches(1.25)))
    footer(s, 13)

# ---- CASE 03 ----
def s14_case03_a():
    s = slide()
    case_band(s, "CASE 03", ["外車の修理費対応の中で評価損の請求漏れが見つかった事例"])
    txt(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(1.4),
        [P("修理費が支払われれば損害賠償は完結する——そう思われがちですが、外車・高額車両・購入間もない車両の場合は、「評価損」という論点が生じることがあります。弁護士が介入したことで、この請求漏れに気づいた事例です。", 17, TEXT_SUB, False, SANS)],
        line_spacing=1.45, space_after=Pt(0))
    num_badge(s, Inches(0.6), Inches(3.65), 1)
    h3(s, Inches(1.3), Inches(3.67), Inches(11), "事案の概要")
    txt(s, Inches(1.3), Inches(4.25), Inches(11.4), Inches(2.5),
        [P("3件目は、自動車修理に関する事例です。この案件では、対象車両が外車であり、しかも購入からまだ1年も経っていない車でした。その車両を修理に出したものの、相手方との調整がなかなかうまく進まず、保険代理店様を通じて弁護士に相談が来たという事案です。", 17, TEXT, False, SANS)],
        line_spacing=1.5, space_after=Pt(0))
    footer(s, 14)

def s15_case03_b():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 03 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 2)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "修理費だけでは足りなかった理由")
    txt(s, Inches(0.6), Inches(1.65), Inches(8.3), Inches(4.0),
        [P("もともとは、修理見積書を作成し、その金額を相手方保険会社に提出して対応してもらう予定でした。しかし、弁護士が介入して内容を確認したところ、修理費の請求だけではなく、評価損も請求できる事例であることが分かりました。外車で、しかも購入後1年未満という条件ですので、修理によって物理的には元どおり使用できるようになったとしても、事故歴が付くことで将来売却時の市場価値が下がり、売れにくくなる可能性があります。つまり、このようなケースでは、単に修理費が支払われれば足りるのではなく、評価損についても併せて請求すべき事例です。", 16, TEXT, False, SANS)],
        line_spacing=1.45, space_after=Pt(0))
    add_picture_cover(s, IMGS["評価損"], Inches(9.2), Inches(1.65), Inches(3.5), Inches(1.8))
    rect(s, Inches(9.2), Inches(3.65), Inches(3.5), Inches(2.5), fill=GREEN_LT)
    txt(s, Inches(9.45), Inches(3.85), Inches(3.05), Inches(2.15),
        [P("評価損が認められやすい条件", 16, GREEN, True, SANS),
         P("外車・高額車両・購入から年数が経っていない車両など、事故後の市場価値の低下が顕著なケース", 16, GREEN_DARK, False, SANS)],
        line_spacing=1.3, space_after=Pt(8))
    footer(s, 15)

def s16_case03_c():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 03 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 3)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "弁護士介入によって見つかった論点")
    txt(s, Inches(0.6), Inches(1.65), Inches(12.1), Inches(2.4),
        [P("この件では、修理見積書を取得し、その見積額を前提に相手方保険会社との交渉があまりうまくいっていないという段階でご相談をいただきました。このタイミングで弁護士が入ったことにより、「評価損の請求がされていない」という点に気づき、追加請求を行うことができました。その意味では、途中段階であっても適切なタイミングで弁護士が介入したことにより、依頼者にとって有利な解決につながった事例といえます。", 17, TEXT, False, SANS)],
        line_spacing=1.45, space_after=Pt(0))
    num_badge(s, Inches(0.6), Inches(4.3), 4)
    h3(s, Inches(1.3), Inches(4.32), Inches(11), "この事例が示す共通課題")
    txt(s, Inches(1.3), Inches(4.9), Inches(11.4), Inches(2.0),
        [P("もっとも、この事例も、これまでの2つの事例と共通するところがあります。すなわち、途中からバトンを渡されて弁護士が追加で入ると、相手方保険会社としても、当初想定していなかった論点への対応を迫られることになり、調整に時間がかかりやすくなります。修理費だけで進むと思っていた案件に、後から評価損の話が入ってくるわけですから、対応が複雑になってしまいます。", 16, TEXT, False, SANS)],
        line_spacing=1.4, space_after=Pt(0))
    footer(s, 16)

def s17_case03_d():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Case 03 / 続き")
    num_badge(s, Inches(0.6), Inches(1.0), 5)
    h3(s, Inches(1.3), Inches(1.02), Inches(11), "初動で見越しておくべきポイント")
    txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.5),
        [P("したがって、この案件についても、やはり最初の段階から弁護士が関与していた方が、全体としてはよりスムーズだったのではないかと思います。特に、外車、高額車両、購入間もない車両などは、修理費だけではなく評価損の問題が生じる可能性がありますので、その点を最初から見越して進めることが重要です。", 16, TEXT, False, SANS)],
        line_spacing=1.4, space_after=Pt(0))
    # 緑見出し
    txt(s, Inches(0.6), Inches(3.15), Inches(12), Inches(0.5),
        [P("評価損の請求を検討すべき主なケース", 20, GREEN_DARK, True, SERIF)],
        line_spacing=1.0, space_after=Pt(0))
    cols = [
        ("外車・輸入車", "事故歴による市場価値低下が国産車より大きくなるケースがある"),
        ("購入から年数が浅い車両", "購入後1〜2年以内は評価損が認められる可能性が特に高い"),
        ("修理費が高額なケース", "修理費が車両価格の一定割合を超える損傷は評価損の根拠になりやすい"),
    ]
    cx = Inches(0.6); cw = Inches(3.9); gap = Inches(0.2)
    for i, (head, body) in enumerate(cols):
        x = Emu(int(cx) + i * (int(cw) + int(gap)))
        rect(s, x, Inches(3.75), cw, Inches(1.85), fill=GREEN_LT)
        txt(s, Emu(int(x)+int(Inches(0.25))), Inches(3.95), Emu(int(cw)-int(Inches(0.5))), Inches(0.5),
            [P(head, 18, GREEN, True, SERIF)], line_spacing=1.1, space_after=Pt(4))
        txt(s, Emu(int(x)+int(Inches(0.25))), Inches(4.5), Emu(int(cw)-int(Inches(0.5))), Inches(1.0),
            [P(body, 16, TEXT, False, SANS)], line_spacing=1.3, space_after=Pt(0))
    # 濃緑 保険代理店様へのポイント
    rect(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(1.25), fill=GREEN_DARK)
    txt(s, Inches(0.95), Inches(5.9), Inches(7.1), Inches(1.0),
        [P("保険代理店様へのポイント", 16, GREEN_MID, True, SANS),
         P("外車・高額車両・購入間もない車両の物損案件は、修理費の見積もり段階から速やかに弁護士にご相談ください。", 16, WHITE, False, SANS)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2, space_after=Pt(2))
    checks = ["CHECK 01  購入から2年以内か", "CHECK 02  外車・輸入車か", "CHECK 03  修理費が高額か"]
    txt(s, Inches(8.3), Inches(5.85), Inches(4.1), Inches(1.05),
        [P(c, 16, WHITE, True, SANS) for c in checks],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=Pt(2))
    footer(s, 17)

# ---- 結論 ----
def s18_summary():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Conclusion")
    heading(s, Inches(0.6), Inches(0.95), Inches(11), "3つの事例に共通していること", size=30)
    h3(s, Inches(0.6), Inches(2.0), Inches(8), "共通する結論", color=GREEN)
    txt(s, Inches(0.6), Inches(2.65), Inches(8.3), Inches(4.0),
        [P("以上の3つの事例に共通しているのは、いずれも「途中から弁護士が入っても解決はできるが、最初から弁護士が入っていた方が、より早く、よりスムーズに進められたはずだ」という点です。治療費打切り後の請求漏れの事例でも、後遺障害診断書の修正事例でも、外車修理における評価損の事例でも、途中からの介入によって救済できた部分はありました。しかし、その一方で、最初から関与していれば避けられた時間的ロスや調整負担があったことも事実です。", 18, TEXT, False, SANS)],
        line_spacing=1.55, space_after=Pt(0))
    add_picture_cover(s, IMGS["家族"], Inches(9.2), Inches(2.0), Inches(3.5), Inches(4.3))
    footer(s, 18)

def s19_table():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Conclusion")
    h3(s, Inches(0.6), Inches(0.95), Inches(11), "3つの事例の比較", size=24, color=GREEN)
    rows = [
        ["事例", "途中介入で生じた課題", "初動介入であれば"],
        ["事例1\n治療費打切り・請求漏れ", "請求後ろ倒し・保険会社の社内調整やり直し・案件長期化", "当初から全請求項目を整理・スムーズな折衝が可能"],
        ["事例2\n後遺障害認定", "一度不認定→やり直し→余分な時間・労力", "認定を見据えた診断書・資料整備で一発認定を狙える"],
        ["事例3\n外車・評価損", "修理費のみで交渉開始→後から評価損追加で複雑化", "評価損を当初から請求項目に組み込み整合的に交渉できる"],
    ]
    rcount = len(rows); ccount = 3
    gtbl = s.shapes.add_table(rcount, ccount, Inches(0.6), Inches(1.7),
                              Inches(12.1), Inches(4.9)).table
    gtbl.columns[0].width = Inches(2.6)
    gtbl.columns[1].width = Inches(4.75)
    gtbl.columns[2].width = Inches(4.75)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
            cell.margin_top = Inches(0.08); cell.margin_bottom = Inches(0.08)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            lines = val.split("\n")
            for li, line in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.line_spacing = 1.2
                r = p.add_run(); r.text = line
                if ri == 0:
                    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE
                elif ci == 0:
                    r.font.size = Pt(17); r.font.bold = (li == 0)
                    r.font.color.rgb = GREEN_DARK if li == 0 else TEXT_SUB
                else:
                    r.font.size = Pt(16); r.font.color.rgb = TEXT
                _set_ea_font(r, SANS if not (ri == 0 or (ci == 0 and li == 0)) else SERIF)
            # セル背景
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = GREEN
            elif ci == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = GREEN_LT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    footer(s, 19)

def s20_benefits():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Conclusion / 続き")
    h3(s, Inches(0.6), Inches(0.95), Inches(11), "最初から弁護士が関与することの利点", color=GREEN, size=24)
    txt(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.6),
        [P("弁護士費用特約が使える案件であれば、「揉めてから弁護士に依頼する」という発想ではなく、「揉める前から弁護士を入れておく」という発想の方が、依頼者にとっても、紹介元にとっても、相手方保険会社にとっても、結果的には合理的であることが多いです。実際、最初から弁護士が入っていれば、請求項目の漏れを防ぎやすくなり、後遺障害認定も見据えた準備がしやすくなり、物損についても評価損のような見落としやすい論点を早期に拾うことができます。", 16, TEXT, False, SANS)],
        line_spacing=1.4, space_after=Pt(0))
    claims = [
        ("請求項目の漏れを防ぎやすくなる", "弁護士は交通事故案件の全請求項目を把握しており、初期段階から漏れのない請求リストを作成できます。依頼者が「言われていないから気にしなかった」という状況を防ぎます。"),
        ("後遺障害認定を見据えた準備がしやすくなる", "認定基準を理解した弁護士が主治医と連携することで、診断書の記載方針・必要な検査・資料収集を初動から最適化できます。"),
        ("物損における評価損などを早期に拾いやすくなる", "修理費以外に評価損・代車費用・格落ち損害など見落とされやすい論点も初期段階から整理し、一体的に請求できます。"),
    ]
    y = Inches(3.3)
    for title, body in claims:
        rect(s, Inches(0.6), y, Inches(12.1), Inches(1.12), fill=GREEN_LT)
        rect(s, Inches(0.85), Emu(int(y)+int(Inches(0.34))), Inches(0.45), Inches(0.45),
             fill=GREEN, shape=MSO_SHAPE.OVAL)
        txt(s, Inches(0.85), Emu(int(y)+int(Inches(0.34))), Inches(0.45), Inches(0.45),
            [P("✓", 18, WHITE, True, SANS)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
        txt(s, Inches(1.55), Emu(int(y)+int(Inches(0.08))), Inches(11.2), Inches(0.98),
            [[(title + "　", 16, GREEN_DARK, True, SERIF)],
             [(body, 16, TEXT, False, SANS)]],
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=Pt(2))
        y = Emu(int(y) + int(Inches(1.25)))
    footer(s, 20)

def s21_message():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Message")
    heading(s, Inches(0.6), Inches(0.95), Inches(11), "初動対応を担う方々へのメッセージ", size=28)
    txt(s, Inches(0.6), Inches(2.0), Inches(8.3), Inches(4.6),
        [P("その意味で、保険代理店として、事故後の初動でお客様と接する立場の方が、「まだ揉めていないから様子を見よう」と考えるのではなく、「弁護士費用特約があるなら、まず弁護士につないでおこう」と判断していただくことには、大きな意味があると考えています。", 16, TEXT, False, SANS),
         P("弁護士費用特約は、あくまでも「使えるなら使う」ものです。特約がある限り、依頼者の自己負担は基本的に生じません。であれば、「早く使えば使うほど、依頼者のメリットが大きい」という論理は、至極自然なものです。", 16, TEXT, False, SANS),
         P("「揉める前に弁護士へ」——この一言を合言葉に、事故後の初動でお客様を適切な方向へ誘導していただければ、依頼者にとって最善の解決への近道になります。当事務所は、そのような連携をご希望の保険代理店様からのご相談を随時お待ちしております。", 16, TEXT, False, SANS)],
        line_spacing=1.4, space_after=Pt(10))
    add_picture_cover(s, IMGS["相談"], Inches(9.2), Inches(2.0), Inches(3.5), Inches(1.7))
    # 特約メリット
    rect(s, Inches(9.2), Inches(3.85), Inches(3.5), Inches(2.95), fill=GREEN_LT)
    txt(s, Inches(9.45), Inches(4.0), Inches(3.05), Inches(0.6),
        [P("弁護士費用特約を活用するメリット", 16, GREEN, True, SANS)],
        line_spacing=1.15, space_after=Pt(4))
    merits = ["依頼者の自己負担なしで弁護士に依頼できる",
              "早期着手で請求漏れを防ぎ、適正補償を確保できる",
              "紹介元（保険代理店様）への顧客からの信頼向上につながる"]
    txt(s, Inches(9.45), Inches(4.6), Inches(3.05), Inches(2.1),
        [[("・", 16, GREEN, True, SANS), (m, 16, GREEN_DARK, False, SANS)] for m in merits],
        line_spacing=1.3, space_after=Pt(8))
    footer(s, 21)

def s22_message2():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Message / 続き")
    # 3サマリーカード
    cards = [
        "入通院慰謝料・交通費・休業損害・営業損害など、全請求項目を当初から整理。後から追加する手間なし。",
        "認定基準を踏まえた診断書・資料整備を初期から最適化。不認定からのやり直し不要。",
        "外車評価損・代車費用・格落ち損害など見落とされがちな論点も、初期段階から一体的に請求。",
    ]
    cx = Inches(0.6); cw = Inches(3.9); gap = Inches(0.2)
    for i, c in enumerate(cards):
        x = Emu(int(cx) + i * (int(cw) + int(gap)))
        rect(s, x, Inches(1.05), cw, Inches(1.95), fill=WHITE, line=GREEN_MID, line_w=Pt(1.5))
        rect(s, x, Inches(1.05), cw, Inches(0.12), fill=GREEN)
        txt(s, Emu(int(x)+int(Inches(0.25))), Inches(1.3), Emu(int(cw)-int(Inches(0.5))), Inches(1.6),
            [P(c, 16, TEXT, False, SANS)], line_spacing=1.4, space_after=Pt(0))
    # 大プルクオート
    rect(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(1.6), fill=GREEN)
    txt(s, Inches(1.0), Inches(3.4), Inches(11.3), Inches(1.4),
        [P("「揉めてから弁護士」ではなく「揉める前から弁護士」——弁護士費用特約があるなら、早期の弁護士介入が依頼者にとっても、紹介元にとっても、最も合理的な選択です。", 20, WHITE, True, SERIF)],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4, space_after=Pt(0))
    # 2カラム情報ボックス
    rect(s, Inches(0.6), Inches(5.2), Inches(5.95), Inches(1.55), fill=GREEN_LT)
    txt(s, Inches(0.85), Inches(5.35), Inches(5.45), Inches(1.3),
        [P("当事務所の対応", 16, GREEN, True, SANS),
         P("保険代理店様からのご紹介・連携を随時受け付けております。まずはお気軽にご相談ください。", 16, GREEN_DARK, False, SANS)],
        line_spacing=1.3, space_after=Pt(4))
    rect(s, Inches(6.75), Inches(5.2), Inches(5.95), Inches(1.55), fill=GREEN_LT)
    txt(s, Inches(7.0), Inches(5.35), Inches(5.45), Inches(1.3),
        [P("アクセス", 16, GREEN, True, SANS),
         P("東銀座駅 徒歩7分 ／ 新橋駅 徒歩9分 ／ 〒104-0061 東京都中央区銀座8丁目17-5", 16, GREEN_DARK, False, SANS)],
        line_spacing=1.3, space_after=Pt(4))
    footer(s, 22)

def s23_contact():
    s = slide()
    # 上：CTA帯（濃緑・背景写真）
    add_picture_cover(s, IMGS["cta"], 0, 0, SW, Inches(2.7))
    # 濃緑半透明オーバーレイの代わりに濃緑帯を左に重ねる（視認性）
    band = rect(s, 0, 0, SW, Inches(2.7), fill=GREEN_DARK)
    band.fill.fore_color.rgb = GREEN_DARK
    # 写真を活かすため帯を消し、左に濃緑グラデ風ボックスのみ置く
    band._element.getparent().remove(band._element)
    leftbox = rect(s, 0, 0, Inches(8.5), Inches(2.7), fill=GREEN_DARK)
    txt(s, Inches(0.6), Inches(0.35), Inches(7.5), Inches(0.4),
        [P("CONTACT", 16, GREEN_MID, True, SANS)], line_spacing=1.0, space_after=Pt(0))
    txt(s, Inches(0.6), Inches(0.75), Inches(7.7), Inches(1.0),
        [P("保険代理店様からのご相談・連携のお問い合わせ", 24, WHITE, True, SERIF)],
        line_spacing=1.2, space_after=Pt(0))
    txt(s, Inches(0.6), Inches(1.75), Inches(7.7), Inches(0.9),
        [P("交通事故案件における弁護士費用特約の活用や、お客様へのご紹介の流れについて、ご不明な点はお気軽にお問い合わせください。初回の相談は無料で承ります。対面・オンライン（Zoom等）どちらでも対応しております。", 16, GREEN_LT, False, SANS)],
        line_spacing=1.3, space_after=Pt(0))
    # バッジ4つ（CTA帯の下端に）
    badges = ["オンライン全国対応", "初回相談無料", "土日祝・夜間対応（要事前相談）", "弁護士費用特約案件 取扱多数"]
    # 下半分：事務所概要 / 本レポートについて
    txt(s, Inches(0.6), Inches(2.85), Inches(11), Inches(0.5),
        [[("弁護士法人 相生綜合法律事務所　", 18, GREEN_DARK, True, SERIF),
          ("AIOI SOGO LAW OFFICE", 16, TEXT_SUB, False, SANS)]],
        line_spacing=1.1, space_after=Pt(0))
    # バッジ
    bx = Inches(0.6)
    for b in badges:
        w = Inches(0.32 + 0.13 * len(b))
        rect(s, bx, Inches(3.4), w, Inches(0.42), fill=GREEN_LT)
        txt(s, bx, Inches(3.4), w, Inches(0.42),
            [P(b, 16, GREEN_DARK, True, SANS)], align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=Pt(0))
        bx = Emu(int(bx) + int(w) + int(Inches(0.15)))
    # 事務所概要（全幅・2カラムのラベル：値）
    txt(s, Inches(0.6), Inches(4.1), Inches(6), Inches(0.35),
        [P("ABOUT", 16, GREEN, True, SANS)], line_spacing=1.0, space_after=Pt(0))
    txt(s, Inches(0.6), Inches(4.45), Inches(6), Inches(0.5),
        [P("事務所概要", 22, GREEN_DARK, True, SERIF)], line_spacing=1.0, space_after=Pt(0))
    rows_left = [
        ("所在地", "〒104-0061 東京都中央区銀座8丁目17-5"),
        ("事務所名", "弁護士法人 相生綜合法律事務所"),
        ("代表弁護士", "福島 駿太（第一東京弁護士会 登録）"),
    ]
    rows_right = [
        ("取扱分野", "交通事故、離婚・相続、労働問題、債務整理、企業法務 ほか"),
        ("対応エリア", "オンラインにて全国対応"),
        ("初回相談", "無料（対面・オンライン対応）"),
    ]
    txt(s, Inches(0.6), Inches(5.2), Inches(6.2), Inches(2.2),
        [[(lbl + "：", 17, GREEN, True, SANS), (val, 17, TEXT, False, SANS)] for lbl, val in rows_left],
        line_spacing=1.35, space_after=Pt(10))
    txt(s, Inches(7.0), Inches(5.2), Inches(5.85), Inches(2.2),
        [[(lbl + "：", 17, GREEN, True, SANS), (val, 17, TEXT, False, SANS)] for lbl, val in rows_right],
        line_spacing=1.35, space_after=Pt(10))

def s24_disclaimer():
    s = slide()
    section_label(s, Inches(0.6), Inches(0.5), "Disclaimer")
    heading(s, Inches(0.6), Inches(0.95), Inches(11), "本レポートについて", size=30)
    txt(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(4.5),
        [P("本レポートに記載されている事例は、実際の案件をもとに一部内容を整理・匿名化したものです。個別の案件の結果を保証するものではありません。", 18, TEXT, False, SANS),
         P("記載内容は2026年4月時点の情報をもとにしており、法令・実務の変更により内容が変わる場合があります。具体的なご相談については、弁護士にご相談ください。", 18, TEXT, False, SANS),
         P("本レポートの著作権は弁護士法人相生綜合法律事務所に帰属します。無断転載・複製はお断りしております。", 18, TEXT, False, SANS)],
        line_spacing=1.55, space_after=Pt(18))
    # 末尾に事務所名
    rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.85), fill=GREEN_LT)
    txt(s, Inches(0.95), Inches(6.0), Inches(11.4), Inches(0.85),
        [[("弁護士法人 相生綜合法律事務所　", 18, GREEN_DARK, True, SERIF),
          ("AIOI SOGO LAW OFFICE", 16, TEXT_SUB, False, SANS)]],
        anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=Pt(0))

# ================================================================ 検査・保存

def verify_min_font():
    """全スライドの全テキストrunが16pt以上か検査。違反があれば例外。"""
    bad = []
    for si, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size is not None and r.font.size < Pt(MIN_PT):
                            bad.append((si, r.text[:20], r.font.size.pt))
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if r.font.size is not None and r.font.size < Pt(MIN_PT):
                                    bad.append((si, r.text[:20], r.font.size.pt))
    if bad:
        for b in bad:
            print("  16pt未満:", b)
        raise SystemExit(f"[NG] 16pt未満のrunが{len(bad)}件あります。")
    print("[OK] 全テキスト16pt以上")

def build_all():
    s01_cover(); s02_intro1(); s03_intro2(); s04_toc()
    s05_case01_a(); s06_case01_b(); s07_case01_c(); s08_case01_d(); s09_case01_e()
    s10_case02_a(); s11_case02_b(); s12_case02_c(); s13_case02_d()
    s14_case03_a(); s15_case03_b(); s16_case03_c(); s17_case03_d()
    s18_summary(); s19_table(); s20_benefits(); s21_message(); s22_message2()
    s23_contact(); s24_disclaimer()

if __name__ == "__main__":
    import sys
    build_all()
    verify_min_font()
    try:
        prs.save(OUT)
    except Exception as e:
        raise SystemExit(f"[NG] 保存失敗（ファイルが開かれている可能性）: {e}")
    print(f"[OK] 保存しました: {OUT}  スライド数={len(prs.slides._sldIdLst)}")
